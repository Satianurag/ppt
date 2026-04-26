"""Template loading, preservation of fixed slides, and cover population.

The three hackathon Slide Master templates each ship with a cover layout and a
thank-you (end) layout. Under C1, slide 1 and slide 15 must remain *visually*
identical to the template (same layout, same master styling). Only the cover
receives four dynamic runtime fields: title, subtitle, presenter name, date.

This module owns everything that is template-shaped and template-sensitive.
"""

from __future__ import annotations

from enum import Enum
import os
import zipfile
from typing import NamedTuple, Optional

from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


class TemplateType(str, Enum):
    AI_BUBBLE = "AI_Bubble"
    UAE_SOLAR = "UAE_Solar"
    ACCENTURE = "Accenture"
    UNKNOWN = "unknown"


class LayoutSlot(NamedTuple):
    """A usable layout in the master, with its declared role."""
    index: int
    name: str
    role: str  # "cover" | "divider" | "content" | "end"


# Per-template layout role map. Only fixed slots (cover, end) are authoritative;
# content slots are a *pool* the scheduler picks from.
_TEMPLATE_LAYOUTS: dict[TemplateType, dict[str, list[int]]] = {
    TemplateType.AI_BUBBLE: {
        "cover": [0],
        "divider": [1],
        "content": [2, 3],   # Blank, Title only
        "end": [4],          # 1_Thank you
    },
    TemplateType.UAE_SOLAR: {
        "cover": [0],
        "divider": [1],
        "content": [2, 3],   # body layouts with 2-4 placeholders
        "end": [4],          # Thank-you text embedded in layout
    },
    TemplateType.ACCENTURE: {
        "cover": [0],
        "divider": [2],
        "content": [3, 4],   # Blank, Title only
        "end": [5],          # Thank You
    },
}


# Cover placeholder slots per template. UAE uses CENTER_TITLE (idx 0) and has no
# subtitle placeholder, so subtitle is rendered as a runtime textbox beneath it.
_COVER_PLACEHOLDERS: dict[TemplateType, dict[str, Optional[int]]] = {
    TemplateType.AI_BUBBLE: {"title": 10, "subtitle": 11},
    TemplateType.UAE_SOLAR: {"title": 0, "subtitle": None},
    TemplateType.ACCENTURE: {"title": 10, "subtitle": 11},
}


def detect_template(prs: Presentation) -> TemplateType:
    """Identify which of the three hackathon templates is loaded."""
    names = [layout.name for layout in prs.slide_masters[0].slide_layouts]
    if "Cover" in names and "Blank" in names:
        return TemplateType.AI_BUBBLE
    if any("0_Title Company" in n for n in names):
        return TemplateType.UAE_SOLAR
    if "1_Cover" in names:
        return TemplateType.ACCENTURE
    return TemplateType.UNKNOWN


def role_layout_indices(template: TemplateType, role: str) -> list[int]:
    """Return the slide-master layout indices available for a role."""
    table = _TEMPLATE_LAYOUTS.get(template)
    if table is None:
        table = _TEMPLATE_LAYOUTS[TemplateType.AI_BUBBLE]
    return list(table.get(role, []))


def get_layout(prs: Presentation, template: TemplateType, role: str, pick: int = 0):
    """Get a slide layout object for a named role (picks the first by default)."""
    indices = role_layout_indices(template, role)
    if not indices:
        raise ValueError(f"No layout registered for role {role!r} in {template.value}")
    idx = indices[min(pick, len(indices) - 1)]
    return prs.slide_masters[0].slide_layouts[idx]


def _drop_all_slides(prs: Presentation) -> None:
    """Delete every existing slide in the presentation (XML-level)."""
    xml_slides = prs.slides._sldIdLst
    rel_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    while len(xml_slides) > 0:
        r_id = xml_slides[0].get(rel_ns)
        if r_id is not None:
            prs.part.drop_rel(r_id)
        del xml_slides[0]


def load_blank_canvas(template_path: str) -> tuple[Presentation, TemplateType]:
    """Load a template and strip all existing demo slides.

    Returns the Presentation with zero slides plus the detected template type.
    The slide master and all layouts are preserved; only slide-level content is
    removed so we can rebuild the 15-slide spine from scratch.
    """
    prs = Presentation(template_path)
    tpl = detect_template(prs)
    _drop_all_slides(prs)
    return prs, tpl


def load_editable_canvas(template_path: str) -> tuple[Presentation, TemplateType]:
    """Create a picture-free 16:9 canvas while preserving template identity.

    The supplied master decks contain media in masters/layouts. Because the final
    output must contain no pictures at all, generation starts from a clean PPTX
    package and later copies only the selected master's theme XML.
    """
    source = Presentation(template_path)
    tpl = detect_template(source)
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    _drop_all_slides(prs)
    return prs, tpl


def apply_theme_from_template(pptx_path: str, template_path: str) -> None:
    """Copy only theme XML from the selected master into the generated deck."""
    theme_xml = _read_theme_xml(template_path)
    if theme_xml is None:
        return

    tmp_path = f"{pptx_path}.theme.tmp"
    with zipfile.ZipFile(pptx_path, "r") as zin:
        with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "ppt/theme/theme1.xml":
                    data = theme_xml
                zout.writestr(item, data)
    os.replace(tmp_path, pptx_path)


def _read_theme_xml(template_path: str) -> bytes | None:
    with zipfile.ZipFile(template_path) as zf:
        theme_names = sorted(
            n for n in zf.namelist()
            if n.startswith("ppt/theme/theme") and n.endswith(".xml")
        )
        if not theme_names:
            return None
        return zf.read(theme_names[0])


def add_native_content_slide(prs: Presentation):
    """Add a blank slide on the clean editable canvas."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    _remove_empty_placeholders(slide)
    return slide


def add_native_cover_slide(
    prs: Presentation,
    template: TemplateType,
    title: str,
    subtitle: Optional[str],
    presenter: str,
    presentation_date: str,
):
    """Editable, picture-free cover inspired by the selected master."""
    slide = add_native_content_slide(prs)
    _add_brand_backdrop(slide, template)
    if template == TemplateType.UAE_SOLAR:
        _add_sunburst(slide, Inches(10.7), Inches(1.0), Inches(1.7))
        title_top = Inches(2.2)
    elif template == TemplateType.ACCENTURE:
        _add_angle_band(slide, Inches(8.0), Inches(0), Inches(5.4), Inches(7.5))
        title_top = Inches(1.7)
    else:
        _add_bubbles(slide)
        title_top = Inches(1.8)

    _add_runtime_textbox(
        slide,
        left=Inches(0.75), top=title_top,
        width=Inches(8.8), height=Inches(1.6),
        text=title,
        font_size=Pt(42),
        bold=True,
    )
    if subtitle:
        _add_runtime_textbox(
            slide,
            left=Inches(0.8), top=Emu(title_top + Inches(1.7)),
            width=Inches(8.6), height=Inches(0.8),
            text=subtitle,
            font_size=Pt(19),
            line_spacing=1.15,
        )
    _add_runtime_textbox(
        slide,
        left=Inches(0.8), top=Inches(6.45),
        width=Inches(5.6), height=Inches(0.55),
        text=f"{presenter}  |  {presentation_date}",
        font_size=Pt(12),
    )
    return slide


def add_agenda_slide(prs: Presentation, template: TemplateType, items: list[str]):
    """Fixed slide 2: editable agenda/index."""
    slide = add_native_content_slide(prs)
    _add_brand_header(slide, template, "Agenda")
    agenda_items = items[:6] or ["Context", "Insights", "Recommendations"]
    top = Inches(1.75)
    for i, item in enumerate(agenda_items):
        y = Emu(top + i * Inches(0.75))
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), y, Inches(0.45), Inches(0.45))
        _theme_fill(badge, _primary_accent(template))
        _shape_text(badge, f"{i + 1}", 12, bold=True, align=PP_ALIGN.CENTER,
                    font_theme=MSO_THEME_COLOR.BACKGROUND_1)
        _add_runtime_textbox(
            slide,
            left=Inches(1.65), top=Emu(y - Inches(0.03)),
            width=Inches(10.4), height=Inches(0.5),
            text=item,
            font_size=Pt(19),
        )
    return slide


def add_summary_slide(prs: Presentation, template: TemplateType, items: list[str]):
    """Fixed slide 14: editable summary/takeaways."""
    slide = add_native_content_slide(prs)
    _add_brand_header(slide, template, "Executive summary")
    takeaways = items[:4] or ["Key findings are consolidated here."]
    gap = Inches(0.25)
    card_w = Emu((Inches(12.0) - gap) / 2)
    card_h = Inches(1.55)
    for i, text in enumerate(takeaways):
        row = i // 2
        col = i % 2
        left = Emu(Inches(0.75) + col * (card_w + gap))
        top = Emu(Inches(1.85) + row * Inches(1.95))
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        _theme_fill(card, _soft_accent(i))
        card.line.fill.background()
        _shape_text(card, text, 16, bold=True, align=PP_ALIGN.CENTER,
                    font_theme=MSO_THEME_COLOR.TEXT_1)
    return slide


def add_native_end_slide(prs: Presentation, template: TemplateType):
    """Editable, picture-free thank-you slide."""
    slide = add_native_content_slide(prs)
    _add_brand_backdrop(slide, template)
    if template == TemplateType.AI_BUBBLE:
        _add_bubbles(slide)
    elif template == TemplateType.UAE_SOLAR:
        _add_sunburst(slide, Inches(10.4), Inches(1.0), Inches(1.6))
    else:
        _add_angle_band(slide, Inches(8.7), Inches(0), Inches(4.8), Inches(7.5))
    _add_runtime_textbox(
        slide,
        left=Inches(1.0), top=Inches(2.65),
        width=Inches(7.8), height=Inches(1.2),
        text="Thank you",
        font_size=Pt(54),
        bold=True,
    )
    _add_runtime_textbox(
        slide,
        left=Inches(1.05), top=Inches(3.95),
        width=Inches(7.8), height=Inches(0.6),
        text="Questions and discussion",
        font_size=Pt(20),
    )
    return slide


def _blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[-1]


def _primary_accent(template: TemplateType) -> MSO_THEME_COLOR:
    if template == TemplateType.UAE_SOLAR:
        return MSO_THEME_COLOR.ACCENT_2
    return MSO_THEME_COLOR.ACCENT_1


def _soft_accent(index: int) -> MSO_THEME_COLOR:
    return [
        MSO_THEME_COLOR.ACCENT_1,
        MSO_THEME_COLOR.ACCENT_2,
        MSO_THEME_COLOR.ACCENT_3,
        MSO_THEME_COLOR.ACCENT_4,
    ][index % 4]


def _theme_fill(shape, color: MSO_THEME_COLOR, transparency: float = 0.0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.theme_color = color
    shape.fill.transparency = transparency


def _shape_text(
    shape,
    text: str,
    size: int,
    *,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_theme: MSO_THEME_COLOR | None = None,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        if font_theme is not None:
            run.font.color.theme_color = font_theme


def _add_brand_header(slide, template: TemplateType, title: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.22))
    _theme_fill(bar, _primary_accent(template))
    bar.line.fill.background()
    _add_runtime_textbox(
        slide,
        left=Inches(0.75), top=Inches(0.48),
        width=Inches(11.8), height=Inches(0.75),
        text=title,
        font_size=Pt(34),
        bold=True,
    )


def _add_brand_backdrop(slide, template: TemplateType) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    _theme_fill(bg, MSO_THEME_COLOR.BACKGROUND_1)
    bg.line.fill.background()
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), Inches(7.5))
    _theme_fill(rail, _primary_accent(template))
    rail.line.fill.background()


def _add_bubbles(slide) -> None:
    specs = [
        (9.0, 0.8, 2.6, 0.22),
        (10.9, 2.7, 1.4, 0.30),
        (8.6, 5.1, 2.0, 0.35),
        (11.5, 5.5, 0.85, 0.15),
    ]
    for x, y, size, transparency in specs:
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        _theme_fill(oval, MSO_THEME_COLOR.ACCENT_1, transparency)
        oval.line.fill.background()


def _add_angle_band(slide, left, top, width, height) -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, left, top, width, height)
    _theme_fill(band, MSO_THEME_COLOR.ACCENT_1, 0.08)
    band.line.fill.background()


def _add_sunburst(slide, cx, cy, radius) -> None:
    sun = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, radius, radius)
    _theme_fill(sun, MSO_THEME_COLOR.ACCENT_2, 0.05)
    sun.line.fill.background()
    center_x = Emu(cx + radius / 2)
    center_y = Emu(cy + radius / 2)
    for angle in range(0, 360, 30):
        import math
        length = Inches(1.15)
        x2 = Emu(center_x + int(math.cos(math.radians(angle)) * length))
        y2 = Emu(center_y + int(math.sin(math.radians(angle)) * length))
        ray = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, center_x, center_y, x2, y2)
        ray.line.color.theme_color = MSO_THEME_COLOR.ACCENT_2
        ray.line.width = Pt(1)


def add_cover_slide(
    prs: Presentation,
    template: TemplateType,
    title: str,
    subtitle: Optional[str],
    presenter: str,
    presentation_date: str,
):
    """Add slide 1 using the template's cover layout and fill the dynamic fields.

    The cover retains the template's master styling entirely — we only write
    into placeholders (title, subtitle if available) and add a textbox at the
    bottom-left for the presenter name and date.
    """
    layout = get_layout(prs, template, "cover")
    slide = prs.slides.add_slide(layout)

    slots = _COVER_PLACEHOLDERS.get(template, {})
    title_idx = slots.get("title")
    subtitle_idx = slots.get("subtitle")

    subtitle_placed = False
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if title_idx is not None and idx == title_idx:
            ph.text = title
        elif subtitle_idx is not None and idx == subtitle_idx and subtitle:
            ph.text = subtitle
            subtitle_placed = True

    if subtitle and not subtitle_placed:
        _add_runtime_textbox(
            slide,
            left=Inches(0.5), top=Inches(4.7),
            width=Inches(12.33), height=Inches(0.8),
            text=subtitle,
            font_size=Pt(20),
        )

    _add_runtime_textbox(
        slide,
        left=Inches(0.4), top=Inches(6.5),
        width=Inches(6.5), height=Inches(0.8),
        text=f"{presenter}\n{presentation_date}",
        font_size=Pt(14),
        line_spacing=1.15,
    )

    return slide


def add_end_slide(prs: Presentation, template: TemplateType):
    """Append slide 15 using the template's thank-you layout, unmodified."""
    layout = get_layout(prs, template, "end")
    return prs.slides.add_slide(layout)


def add_content_slide(prs: Presentation, template: TemplateType, pick: int = 0):
    """Add a body slide using a neutral content layout from the master.

    The layout is a blank-ish canvas — the slide's visual identity is produced
    by the layout renderer (which draws the geometry), while the master still
    provides theme colors, fonts, and the background chrome.

    Any title/body placeholders inherited from the layout are removed so empty
    "Click to add Title" / "Click to add Text" prompts do not ghost through the
    rendered geometry.
    """
    layout = get_layout(prs, template, "content", pick=pick)
    slide = prs.slides.add_slide(layout)
    _remove_empty_placeholders(slide)
    return slide


def _remove_empty_placeholders(slide) -> None:
    """Delete every placeholder shape from a slide.

    The layout's master chrome (logos, background, decorative shapes) is part of
    the slide layout XML, not the slide itself — removing slide-level
    placeholders does not touch the master chrome.
    """
    spTree = slide.shapes._spTree
    for ph in list(slide.placeholders):
        spTree.remove(ph._element)


def _add_runtime_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_size,
    bold: bool = False,
    line_spacing: float = 1.0,
) -> None:
    """Add a plain textbox with zero local styling overrides.

    Deliberately does NOT set ``font.name`` or ``font.color`` — both cascade
    from the slide master (C3). ``font.size`` is the only run-level override we
    allow, because master-inherited sizes would otherwise render the cover's
    bottom-left block too large.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = font_size
            if bold:
                run.font.bold = True
