"""14-layout catalog — prebuilt layout components for body slides.

Each layout is a *geometric class*. Two adjacent slides may not share a class
(constraint C2). Layouts render with python-pptx shapes that inherit fonts and
colors from the Slide Master (constraint C3): no runtime srgbClr, no font.name.

Shape fill uses ``MSO_THEME_COLOR`` references so colors cascade from the
master theme. Only font sizes are locally overridden (master-inherited sizes
would wreck the infographic geometry).

Layout classes (14):
    1.  bullet_text        — title + bullet list
    2.  kpi_big_numbers    — title + 2–4 huge KPI tiles
    3.  four_column_icons  — title + 4 icon tiles with captions
    4.  three_column_grid  — title + 3 column cards
    5.  two_column_compare — title + left/right comparison cards
    6.  chart_focused      — title + full-width chart
    7.  chart_with_bullets — title + chart (left) + bullet list (right)
    8.  table              — title + data table
    9.  timeline           — title + horizontal event timeline
    10. process            — title + numbered step chevrons
    11. pyramid            — title + inverted pyramid tiers
    12. funnel             — title + 3–5 funnel stages
    13. quote_emphasis     — large centered quote
    14. section_divider    — title-only / big typography divider
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import Callable, Optional

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from step2.slide_plan_models import ChartType
from step3.content_models import SlideContent


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TITLE_TOP = Inches(0.5)
TITLE_HEIGHT = Inches(0.9)
TITLE_LEFT = Inches(0.5)
TITLE_WIDTH = Inches(12.333)

BODY_TOP = Inches(1.6)
BODY_LEFT = Inches(0.5)
BODY_WIDTH = Inches(12.333)
BODY_HEIGHT = Inches(5.2)

KEY_FOOTER_TOP = Inches(6.85)
KEY_FOOTER_HEIGHT = Inches(0.4)


# ── Shared primitives ────────────────────────────────────────────────

def _add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float = 1.15,
) -> None:
    """Add a textbox that inherits master font name & color (C3).

    Only size/bold/alignment/anchor/line-spacing are locally controlled — enough
    for geometric layout, but nothing that overrides the master theme.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    lines = text.split("\n") if "\n" in text else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            if bold:
                run.font.bold = True


def _add_title(slide, title: str) -> None:
    _add_text(
        slide, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        title, size=32, bold=True, anchor=MSO_ANCHOR.MIDDLE,
    )


def _add_key_footer(slide, text: str) -> None:
    if not text:
        return
    _add_text(
        slide, BODY_LEFT, KEY_FOOTER_TOP, BODY_WIDTH, KEY_FOOTER_HEIGHT,
        text, size=12, anchor=MSO_ANCHOR.MIDDLE,
    )


_PRIMARY_ACCENT: MSO_THEME_COLOR = MSO_THEME_COLOR.ACCENT_1


def set_primary_accent(theme: MSO_THEME_COLOR) -> None:
    """Set the master theme color used for every primary tile in the deck.

    Templates do not all declare ``accent1`` as their dark/saturated brand
    color — UAE_Solar's ``accent1`` is a near-white sage. The deck builder
    resolves the darkest accent per template at load time and registers it
    here so KPI tiles, cards, and timeline nodes always render with enough
    contrast for white text.
    """
    global _PRIMARY_ACCENT
    _PRIMARY_ACCENT = theme


def _accent_fill(shape, theme: Optional[MSO_THEME_COLOR] = None) -> None:
    """Fill a shape with a master theme color — never a hardcoded RGB."""
    shape.fill.solid()
    shape.fill.fore_color.theme_color = theme if theme is not None else _PRIMARY_ACCENT
    shape.line.fill.background()


def _derive_compare_headers(content: SlideContent) -> tuple[str, str]:
    """Derive meaningful column headers for two-column compare layouts.

    Priority order:
    1. Subtitle contains "vs" / "versus" — split on that.
    2. Title contains "vs" / "versus" — split on that.
    3. First two bullets each start with a label followed by a colon — use those labels.
    4. Fallback: use the first ~3 words of each half's first bullet.
    """
    import re

    for text in (content.subtitle, content.title):
        if not text:
            continue
        for sep in ("versus", " vs. ", " vs ", " vs.", " v. "):
            low = text.lower()
            idx = low.find(sep)
            if idx != -1:
                left = text[:idx].strip()
                right = text[idx + len(sep):].strip()
                if left and right:
                    return left, right

    bullets = _bullet_texts(content)
    if len(bullets) >= 2:
        colon_pat = re.compile(r"^([^:]{2,30}):\s")
        m0 = colon_pat.match(bullets[0])
        m1 = colon_pat.match(bullets[1])
        if m0 and m1:
            return m0.group(1).strip(), m1.group(1).strip()

    if len(bullets) >= 2:
        mid = len(bullets) // 2
        left_words = bullets[0].split()[:4]
        right_words = bullets[mid].split()[:4]
        if left_words and right_words:
            return " ".join(left_words), " ".join(right_words)

    return "Key Strengths", "Key Challenges"


def _bullet_texts(slide_content: SlideContent) -> list[str]:
    """Flatten the Designer output to a list of bullet strings."""
    bullets: list[str] = []
    if slide_content.bullets:
        bullets.extend(b.text for b in slide_content.bullets)
    for kp in slide_content.key_points:
        bullets.extend(kp.bullet_form)
    # Light dedupe, preserve order
    seen: set[str] = set()
    unique: list[str] = []
    for b in bullets:
        if b not in seen:
            seen.add(b)
            unique.append(b)
    return unique


# ── Layout 1: bullet_text ────────────────────────────────────────────

def bullet_text(slide, content: SlideContent) -> None:
    _add_title(slide, content.title)
    bullets = _bullet_texts(content)[:10] or [content.key_message]
    size = 22 if len(bullets) <= 4 else (18 if len(bullets) <= 7 else 14)
    text = "\n".join(f"• {b}" for b in bullets)
    _add_text(slide, BODY_LEFT, BODY_TOP, BODY_WIDTH, BODY_HEIGHT,
              text, size=size, line_spacing=1.35)
    _add_key_footer(slide, content.key_message)


# ── Layout 2: kpi_big_numbers ────────────────────────────────────────

_KPI_VALUE = r"\$?[0-9][\d,.]*\s*(?:%|bn|tn|million|billion|trillion|[KMBT])?"


def _extract_kpi_from_text(text: str) -> Optional[tuple[str, str]]:
    """Return (value, label) if text starts with a self-contained metric.

    A metric is digits optionally prefixed by $, optionally suffixed by %, a
    unit (K/M/B/T/bn/tn) or a full word (million/billion/trillion). A word
    boundary is required between the value and the label so we never chew off
    the first letter of the following word (`326 acquisitions` must not become
    `326 a / cquisitions`).
    """
    import re
    pattern = rf"\s*({_KPI_VALUE})\s+[:—\-–]?\s*(.+)"
    m = re.match(pattern, text, flags=re.IGNORECASE)
    if not m:
        return None
    value = m.group(1).strip()
    label = m.group(2).strip()
    if not value or not label or len(value) > 12:
        return None
    return value, label


def kpi_big_numbers(slide, content: SlideContent) -> None:
    _add_title(slide, content.title)

    kpis: list[tuple[str, str]] = []
    for b in _bullet_texts(content):
        parsed = _extract_kpi_from_text(b)
        if parsed:
            kpis.append(parsed)
    if not kpis:
        # Fallback: synthesise tiles from bullets
        kpis = [(f"{i + 1}", b[:40]) for i, b in enumerate(_bullet_texts(content)[:4])]
    kpis = kpis[:4] or [("—", content.key_message[:60])]

    n = len(kpis)
    gap = Inches(0.3)
    tile_w = Emu(int((BODY_WIDTH - gap * (n - 1)) / n))
    tile_h = Inches(3.4)
    top = Inches(2.2)

    for i, (value, label) in enumerate(kpis):
        left = Emu(BODY_LEFT + i * (tile_w + gap))
        tile = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, tile_w, tile_h)
        _accent_fill(tile)
        tile.text_frame.margin_left = Inches(0.2)
        tile.text_frame.margin_right = Inches(0.2)
        tile.text_frame.margin_top = Inches(0.3)
        tile.text_frame.margin_bottom = Inches(0.3)

        tf = tile.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.text = value
        p1.alignment = PP_ALIGN.CENTER
        for run in p1.runs:
            run.font.size = Pt(54)
            run.font.bold = True
        p2 = tf.add_paragraph()
        p2.text = label
        p2.alignment = PP_ALIGN.CENTER
        for run in p2.runs:
            run.font.size = Pt(14)

    _add_key_footer(slide, content.key_message)


# ── Layout 3: four_column_icons ──────────────────────────────────────

def four_column_icons(slide, content: SlideContent) -> None:
    _add_title(slide, content.title)

    bullets = _bullet_texts(content)[:4] or [content.key_message]
    bullets = (bullets * 4)[:4]

    icons = ["★", "◆", "●", "▲"]
    gap = Inches(0.3)
    col_w = Emu(int((BODY_WIDTH - gap * 3) / 4))
    icon_h = Inches(1.4)
    text_h = Inches(2.0)
    top = Inches(2.0)

    for i, text in enumerate(bullets):
        left = Emu(BODY_LEFT + i * (col_w + gap))
        icon_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, icon_h, icon_h)
        _accent_fill(icon_box)
        icon_box.text_frame.text = icons[i]
        p = icon_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        icon_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        for run in p.runs:
            run.font.size = Pt(48)
            run.font.bold = True

        _add_text(
            slide,
            left, top + icon_h + Inches(0.2), col_w, text_h,
            text, size=14, align=PP_ALIGN.CENTER, line_spacing=1.25,
        )

    _add_key_footer(slide, content.key_message)


# ── Layout 4: three_column_grid ──────────────────────────────────────

def three_column_grid(slide, content: SlideContent) -> None:
    _add_title(slide, content.title)

    bullets = _bullet_texts(content)
    cards = bullets[:3] if len(bullets) >= 3 else (bullets + [content.key_message] * 3)[:3]

    gap = Inches(0.3)
    col_w = Emu(int((BODY_WIDTH - gap * 2) / 3))
    card_h = Inches(4.4)
    top = Inches(1.9)

    for i, text in enumerate(cards):
        left = Emu(BODY_LEFT + i * (col_w + gap))
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, col_w, card_h)
        _accent_fill(card)
        card.text_frame.margin_left = Inches(0.3)
        card.text_frame.margin_right = Inches(0.3)
        card.text_frame.margin_top = Inches(0.3)
        card.text_frame.margin_bottom = Inches(0.3)
        card.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        card.text_frame.word_wrap = True

        tf = card.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = f"{i + 1}"
        p1.alignment = PP_ALIGN.CENTER
        for run in p1.runs:
            run.font.size = Pt(40)
            run.font.bold = True
        p2 = tf.add_paragraph()
        p2.text = text
        p2.alignment = PP_ALIGN.CENTER
        for run in p2.runs:
            run.font.size = Pt(16)

    _add_key_footer(slide, content.key_message)


# ── Layout 5: two_column_compare ─────────────────────────────────────

def two_column_compare(slide, content: SlideContent) -> None:
    _add_title(slide, content.title)
    bullets = _bullet_texts(content)

    left_title, right_title = _derive_compare_headers(content)

    mid = len(bullets) // 2 or 1
    left_bullets = bullets[:mid] or [content.key_message]
    right_bullets = bullets[mid:] or [content.key_message]

    gap = Inches(0.3)
    col_w = Emu(int((BODY_WIDTH - gap) / 2))
    top = Inches(1.9)
    card_h = Inches(4.4)

    for i, (ctitle, citems) in enumerate(
        [
            (left_title, left_bullets),
            (right_title, right_bullets),
        ]
    ):
        left = Emu(BODY_LEFT + i * (col_w + gap))
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_w, card_h)
        _accent_fill(card)
        tf = card.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.3)
        tf.margin_bottom = Inches(0.3)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.word_wrap = True
        tf.clear()

        p1 = tf.paragraphs[0]
        p1.text = ctitle
        p1.alignment = PP_ALIGN.CENTER
        for run in p1.runs:
            run.font.size = Pt(24)
            run.font.bold = True

        for item in citems[:5]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            for run in p.runs:
                run.font.size = Pt(14)
            p.line_spacing = 1.3

    _add_key_footer(slide, content.key_message)


# ── Layouts 6 & 7: chart_focused, chart_with_bullets ─────────────────

_CHART_XL_TYPE = {
    ChartType.BAR: XL_CHART_TYPE.COLUMN_CLUSTERED,
    ChartType.HORIZONTAL_BAR: XL_CHART_TYPE.BAR_CLUSTERED,
    ChartType.GROUPED_BAR: XL_CHART_TYPE.COLUMN_CLUSTERED,
    ChartType.LINE: XL_CHART_TYPE.LINE_MARKERS,
    ChartType.PIE: XL_CHART_TYPE.PIE,
    ChartType.DONUT: XL_CHART_TYPE.DOUGHNUT,
}


def _humanize_large_number_format(fmt: str, max_val: float) -> str:
    """Convert a number format to use B/M/K suffixes for large values."""
    abs_max = abs(max_val) if max_val else 0
    is_currency = "$" in fmt

    if abs_max >= 1_000_000_000:
        if is_currency:
            return '$#,##0.0,,"B"'
        return '#,##0.0,,"B"'
    elif abs_max >= 1_000_000:
        if is_currency:
            return '$#,##0.0,"M"'
        return '#,##0.0,"M"'
    elif abs_max >= 10_000:
        if is_currency:
            return '$#,##0,"K"'
        return '#,##0,"K"'
    return fmt


def _add_chart(slide, left, top, width, height, data) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = data.categories
    for s in data.series:
        chart_data.add_series(s.get("name", ""), s.get("values", []))

    xl_type = _CHART_XL_TYPE.get(data.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    graphic = slide.shapes.add_chart(xl_type, left, top, width, height, chart_data)
    chart = graphic.chart
    chart.has_title = False
    chart.has_legend = data.show_legend
    if data.show_data_labels:
        plot = chart.plots[0]
        plot.has_data_labels = True

    num_fmt = getattr(data, "number_format", "General")
    if num_fmt and num_fmt != "General":
        all_vals = []
        for s in data.series:
            all_vals.extend(v for v in s.get("values", []) if isinstance(v, (int, float)))
        max_val = max(all_vals) if all_vals else 0
        num_fmt = _humanize_large_number_format(num_fmt, max_val)
        try:
            chart.value_axis.tick_labels.number_format = num_fmt
            chart.value_axis.tick_labels.number_format_is_linked = False
        except (ValueError, AttributeError):
            pass


def chart_focused(slide, content: SlideContent) -> None:
    _add_title(slide, content.title)
    if content.chart_data is not None:
        _add_chart(
            slide,
            BODY_LEFT, BODY_TOP, BODY_WIDTH, Inches(5.0),
            content.chart_data,
        )
    else:
        bullet_text(slide, content)
        return
    _add_key_footer(slide, content.key_message)


def chart_with_bullets(slide, content: SlideContent) -> None:
    _add_title(slide, content.title)
    if content.chart_data is None:
        bullet_text(slide, content)
        return

    chart_w = Inches(7.5)
    gap = Inches(0.3)
    right_left = Emu(BODY_LEFT + chart_w + gap)
    right_w = Emu(BODY_WIDTH - chart_w - gap)

    _add_chart(slide, BODY_LEFT, BODY_TOP, chart_w, Inches(5.0), content.chart_data)

    bullets = _bullet_texts(content)[:6] or [content.key_message]
    text = "\n".join(f"• {b}" for b in bullets)
    _add_text(slide, right_left, BODY_TOP, right_w, Inches(5.0),
              text, size=14, line_spacing=1.35)

    _add_key_footer(slide, content.key_message)


# ── Layout 8: table ──────────────────────────────────────────────────

def table_layout(slide, content: SlideContent) -> None:
    _add_title(slide, content.title)
    data = content.table_data
    if data is None:
        bullet_text(slide, content)
        return

    rows = len(data.rows) + 1  # header row
    cols = len(data.headers)
    top = Inches(1.8)
    height = Inches(min(5.0, 0.5 + 0.45 * rows))
    shape = slide.shapes.add_table(rows, cols, BODY_LEFT, top, BODY_WIDTH, height)
    tbl = shape.table

    for j, h in enumerate(data.headers):
        cell = tbl.cell(0, j)
        cell.text = str(h)
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(13)

    for i, row in enumerate(data.rows, start=1):
        for j in range(cols):
            cell = tbl.cell(i, j)
            cell.text = str(row[j]) if j < len(row) else ""
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12)

    _add_key_footer(slide, content.key_message)


# ── Layout 9: timeline ───────────────────────────────────────────────

def timeline(slide, content: SlideContent) -> None:
    _add_title(slide, content.title)
    bullets = _bullet_texts(content)[:6] or [content.key_message]
    n = len(bullets)

    track_y = Inches(4.1)
    track_h = Inches(0.1)
    _accent_fill(
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, BODY_LEFT, track_y, BODY_WIDTH, track_h)
    )

    node_r = Inches(0.45)
    gap = Emu(int((BODY_WIDTH - node_r * n) / max(n - 1, 1)))
    text_h = Inches(1.5)

    for i, b in enumerate(bullets):
        cx = Emu(BODY_LEFT + i * (node_r + gap))
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, track_y - Inches(0.18), node_r, node_r)
        _accent_fill(node)
        node.text_frame.text = str(i + 1)
        p = node.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        node.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        for run in p.runs:
            run.font.bold = True
            run.font.size = Pt(16)

        cap_top = track_y + Inches(0.7) if i % 2 == 0 else Inches(2.2)
        _add_text(
            slide,
            Emu(cx - Inches(0.8)), cap_top, Inches(2.1), text_h,
            b, size=12, align=PP_ALIGN.CENTER, line_spacing=1.2,
        )
    _add_key_footer(slide, content.key_message)


# ── Layout 10: process (numbered chevrons) ───────────────────────────

def process(slide, content: SlideContent) -> None:
    _add_title(slide, content.title)
    bullets = _bullet_texts(content)[:5] or [content.key_message]
    n = len(bullets)

    gap = Inches(0.15)
    chev_w = Emu(int((BODY_WIDTH - gap * (n - 1)) / n))
    chev_h = Inches(1.4)
    top = Inches(2.3)

    for i, b in enumerate(bullets):
        left = Emu(BODY_LEFT + i * (chev_w + gap))
        chev = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, chev_w, chev_h)
        _accent_fill(chev)
        tf = chev.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"{i + 1}. {b}"
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.bold = True

    bullets_text = _bullet_texts(content)[n : n + 6]
    if bullets_text:
        text = "\n".join(f"• {b}" for b in bullets_text)
        _add_text(
            slide, BODY_LEFT, Inches(4.0), BODY_WIDTH, Inches(2.5),
            text, size=14, line_spacing=1.3,
        )

    _add_key_footer(slide, content.key_message)


# ── Layout 11: pyramid ───────────────────────────────────────────────

def pyramid(slide, content: SlideContent) -> None:
    _add_title(slide, content.title)
    levels = _bullet_texts(content)[:4] or [content.key_message]
    n = len(levels)
    base_w = Inches(9.0)
    level_h = Inches(0.9)
    gap = Inches(0.1)
    start_top = Inches(1.9)
    center_x = Emu(BODY_LEFT + BODY_WIDTH // 2)

    for i, text in enumerate(levels):
        w = Emu(int(base_w * (1 - i / (n + 1))))
        left = Emu(center_x - w // 2)
        top = Emu(start_top + i * (level_h + gap))
        shape = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, left, top, w, level_h)
        _accent_fill(shape)
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.bold = True

    _add_key_footer(slide, content.key_message)


# ── Layout 12: funnel ────────────────────────────────────────────────

def funnel(slide, content: SlideContent) -> None:
    _add_title(slide, content.title)
    stages = _bullet_texts(content)[:5] or [content.key_message]
    n = len(stages)
    max_w = Inches(10.0)
    min_w = Inches(3.0)
    gap = Inches(0.15)
    stage_h = Inches(0.75)
    start_top = Inches(2.0)
    center_x = Emu(BODY_LEFT + BODY_WIDTH // 2)

    for i, text in enumerate(stages):
        w = Emu(int(max_w - (max_w - min_w) * (i / max(n - 1, 1))))
        left = Emu(center_x - w // 2)
        top = Emu(start_top + i * (stage_h + gap))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, stage_h)
        _accent_fill(shape)
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.bold = True

    _add_key_footer(slide, content.key_message)


# ── Layout 13: quote_emphasis ────────────────────────────────────────

def quote_emphasis(slide, content: SlideContent) -> None:
    _add_title(slide, content.title)
    quote = content.key_message or (content.bullets[0].text if content.bullets else content.title)
    _add_text(
        slide, Inches(1.5), Inches(2.3), Inches(10.33), Inches(3.5),
        f"“{quote}”", size=32, bold=True, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.35,
    )
    attribution = content.subtitle or "Source"
    _add_text(
        slide, BODY_LEFT, Inches(6.1), BODY_WIDTH, Inches(0.5),
        f"— {attribution}", size=16, align=PP_ALIGN.CENTER,
    )


# ── Layout 14: section_divider ───────────────────────────────────────

def section_divider(slide, content: SlideContent) -> None:
    bullets = _bullet_texts(content)
    has_items = bool(bullets)

    title_top = Inches(2.0) if has_items else Inches(2.8)
    _add_text(
        slide, BODY_LEFT, title_top, BODY_WIDTH, Inches(1.5),
        content.title, size=54, bold=True, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if content.subtitle or content.key_message:
        _add_text(
            slide, BODY_LEFT, Inches(3.5) if has_items else Inches(4.8),
            BODY_WIDTH, Inches(0.8),
            content.subtitle or content.key_message,
            size=22, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

    if has_items:
        items_text = "\n".join(f"•  {b}" for b in bullets[:8])
        _add_text(
            slide, Inches(2.5), Inches(4.4), Inches(8.33), Inches(3.0),
            items_text, size=16, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, line_spacing=1.5,
        )


# ── Catalog ──────────────────────────────────────────────────────────

@dataclass(frozen=True)
class LayoutEntry:
    name: str
    klass: str
    render: Callable[[object, SlideContent], None]


CATALOG: dict[str, LayoutEntry] = {
    "bullet_text":        LayoutEntry("bullet_text",        "bullets",   bullet_text),
    "kpi_big_numbers":    LayoutEntry("kpi_big_numbers",    "kpi",       kpi_big_numbers),
    "four_column_icons":  LayoutEntry("four_column_icons",  "grid-4",    four_column_icons),
    "three_column_grid":  LayoutEntry("three_column_grid",  "grid-3",    three_column_grid),
    "two_column_compare": LayoutEntry("two_column_compare", "compare",   two_column_compare),
    "chart_focused":      LayoutEntry("chart_focused",      "chart",     chart_focused),
    "chart_with_bullets": LayoutEntry("chart_with_bullets", "chart+text", chart_with_bullets),
    "table":              LayoutEntry("table",              "table",     table_layout),
    "timeline":           LayoutEntry("timeline",           "timeline",  timeline),
    "process":            LayoutEntry("process",            "process",   process),
    "pyramid":            LayoutEntry("pyramid",            "pyramid",   pyramid),
    "funnel":             LayoutEntry("funnel",             "funnel",    funnel),
    "quote_emphasis":     LayoutEntry("quote_emphasis",     "quote",     quote_emphasis),
    "section_divider":    LayoutEntry("section_divider",    "divider",   section_divider),
}
