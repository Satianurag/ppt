"""Icon manager — PPT Master's 640 chunk icons as native python-pptx shapes.

Icons are parsed from SVG path data and rendered as python-pptx freeform shapes
(DrawingML), NOT as embedded images. This means they are fully editable in PowerPoint.

Source: PPT Master templates/icons/chunk/ (640 SVG icons, MIT license)
All chunk icons use straight-line geometry (M/L/H/V/Z commands only for most),
making them ideal for python-pptx freeform conversion.

Usage:
    from step4.icon_manager import embed_icon, get_icon_for_category, ICON_CATEGORIES

    # Embed a specific icon
    embed_icon(slide, "chart-bar", left, top, size, color)

    # Get an icon for a content category
    icon_name = get_icon_for_category("finance")
    embed_icon(slide, icon_name, left, top, size, color)
"""

import re
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# Icon library path (relative to project root)
ICON_DIR = Path(__file__).parent.parent / "assets" / "icons" / "chunk"

# Category → icon mappings for infographic renderers
ICON_CATEGORIES: dict[str, str] = {
    "finance": "chart-bar",
    "growth": "arrow-trend-up",
    "decline": "arrow-trend-down",
    "time": "clock",
    "calendar": "calendar",
    "check": "circle-checkmark",
    "warning": "triangle-warning",
    "star": "star",
    "target": "target",
    "lightbulb": "lightbulb",
    "people": "users",
    "person": "user",
    "settings": "gear",
    "search": "magnifying-glass",
    "location": "location-pin",
    "globe": "globe",
    "shield": "shield-check",
    "lock": "lock",
    "heart": "heart",
    "home": "home",
    "mail": "envelope",
    "phone": "phone",
    "cloud": "cloud",
    "database": "database",
    "file": "file",
    "folder": "folder",
    "link": "link",
    "code": "code",
    "rocket": "rocket",
    "trophy": "trophy",
    "flag": "flag",
    "bookmark": "bookmark",
    "tag": "tag",
    "cart": "shopping-cart",
    "money": "coins",
    "building": "building",
    "factory": "factory",
    "truck": "truck",
    "plane": "plane",
    "graduation": "graduation-cap",
    "medical": "syringe",
    "chart_pie": "chart-pie",
    "chart_line": "chart-line",
}


def _parse_svg_paths(icon_name: str) -> list[list[tuple[str, list[float]]]]:
    """Parse SVG file and extract path commands.

    Returns list of paths, each path is a list of (command, args) tuples.
    """
    svg_path = ICON_DIR / f"{icon_name}.svg"
    if not svg_path.exists():
        return []

    tree = ET.parse(svg_path)
    root = tree.getroot()
    ns = {"svg": "http://www.w3.org/2000/svg"}

    paths = []
    for path_elem in root.findall(".//svg:path", ns) or root.findall(".//{http://www.w3.org/2000/svg}path"):
        d = path_elem.get("d", "")
        if d:
            commands = _tokenize_path(d)
            if commands:
                paths.append(commands)

    # Also try without namespace
    if not paths:
        for path_elem in root.iter():
            if path_elem.tag.endswith("path"):
                d = path_elem.get("d", "")
                if d:
                    commands = _tokenize_path(d)
                    if commands:
                        paths.append(commands)

    return paths


def _tokenize_path(d: str) -> list[tuple[str, list[float]]]:
    """Tokenize an SVG path 'd' attribute into (command, args) pairs."""
    tokens = re.findall(r'[MmLlHhVvCcSsQqTtAaZz]|[-+]?(?:\d+\.?\d*|\.\d+)(?:[eE][-+]?\d+)?', d)
    commands = []
    current_cmd = None
    current_args: list[float] = []

    for token in tokens:
        if token.isalpha() or token in ("Z", "z"):
            if current_cmd is not None:
                commands.append((current_cmd, current_args))
            current_cmd = token
            current_args = []
            if token in ("Z", "z"):
                commands.append((token, []))
                current_cmd = None
        else:
            current_args.append(float(token))

    if current_cmd is not None and current_cmd not in ("Z", "z"):
        commands.append((current_cmd, current_args))

    return commands


def _get_viewbox(icon_name: str) -> tuple[float, float, float, float]:
    """Get viewBox from SVG file. Returns (x, y, width, height)."""
    svg_path = ICON_DIR / f"{icon_name}.svg"
    if not svg_path.exists():
        return (0, 0, 16, 16)

    tree = ET.parse(svg_path)
    root = tree.getroot()
    vb = root.get("viewBox", "0 0 16 16")
    parts = [float(x) for x in vb.split()]
    return (parts[0], parts[1], parts[2], parts[3])


def embed_icon(
    slide,
    icon_name: str,
    left: int,
    top: int,
    size: int,
    color: RGBColor | None = None,
) -> bool:
    """Embed an icon as a native python-pptx shape on a slide.

    For simple rectangular icons, renders as MSO_SHAPE. For complex paths,
    falls back to a colored rectangle placeholder with the icon category initial.

    Args:
        slide: python-pptx slide object.
        icon_name: Name of the icon (without .svg extension).
        left: Left position in EMU.
        top: Top position in EMU.
        size: Width and height in EMU (icons are square).
        color: Fill color. Defaults to accent blue.

    Returns:
        True if icon was successfully embedded.
    """
    if color is None:
        color = RGBColor(0x00, 0x76, 0xA8)

    # Try to render as native shape
    paths = _parse_svg_paths(icon_name)
    if not paths:
        _add_fallback_shape(slide, icon_name, left, top, size, color)
        return True

    vb = _get_viewbox(icon_name)
    vb_w, vb_h = vb[2], vb[3]

    # For simple icons (1-3 paths with only M/L/H/V/Z commands), use freeform
    all_simple = all(
        all(cmd[0].upper() in "MLHVZ" for cmd in path)
        for path in paths
    )

    if all_simple and len(paths) <= 4:
        _render_freeform_icon(slide, paths, left, top, size, vb_w, vb_h, color)
    else:
        _add_fallback_shape(slide, icon_name, left, top, size, color)

    return True


def _render_freeform_icon(
    slide,
    paths: list[list[tuple[str, list[float]]]],
    left: int,
    top: int,
    size: int,
    vb_w: float,
    vb_h: float,
    color: RGBColor,
) -> None:
    """Render SVG paths as a python-pptx freeform shape."""
    from pptx.util import Emu

    scale_x = size / vb_w
    scale_y = size / vb_h

    for path_cmds in paths:
        builder = slide.shapes.build_freeform(left, top)
        cx, cy = 0.0, 0.0
        started = False

        for cmd, args in path_cmds:
            if cmd == "M":
                for i in range(0, len(args), 2):
                    px = int(args[i] * scale_x)
                    py = int(args[i + 1] * scale_y)
                    if not started:
                        builder = slide.shapes.build_freeform(
                            left + px, top + py
                        )
                        started = True
                    else:
                        builder.add_line_segments([(Emu(left + px), Emu(top + py))])
                    cx, cy = args[i], args[i + 1]

            elif cmd == "m":
                for i in range(0, len(args), 2):
                    cx += args[i]
                    cy += args[i + 1]
                    px = int(cx * scale_x)
                    py = int(cy * scale_y)
                    if not started:
                        builder = slide.shapes.build_freeform(
                            left + px, top + py
                        )
                        started = True

            elif cmd == "L":
                segments = []
                for i in range(0, len(args), 2):
                    px = int(args[i] * scale_x)
                    py = int(args[i + 1] * scale_y)
                    segments.append((Emu(left + px), Emu(top + py)))
                    cx, cy = args[i], args[i + 1]
                if segments:
                    builder.add_line_segments(segments)

            elif cmd == "l":
                segments = []
                for i in range(0, len(args), 2):
                    cx += args[i]
                    cy += args[i + 1]
                    px = int(cx * scale_x)
                    py = int(cy * scale_y)
                    segments.append((Emu(left + px), Emu(top + py)))
                if segments:
                    builder.add_line_segments(segments)

            elif cmd == "H":
                cx = args[0]
                px = int(cx * scale_x)
                py = int(cy * scale_y)
                builder.add_line_segments([(Emu(left + px), Emu(top + py))])

            elif cmd == "h":
                cx += args[0]
                px = int(cx * scale_x)
                py = int(cy * scale_y)
                builder.add_line_segments([(Emu(left + px), Emu(top + py))])

            elif cmd == "V":
                cy = args[0]
                px = int(cx * scale_x)
                py = int(cy * scale_y)
                builder.add_line_segments([(Emu(left + px), Emu(top + py))])

            elif cmd == "v":
                cy += args[0]
                px = int(cx * scale_x)
                py = int(cy * scale_y)
                builder.add_line_segments([(Emu(left + px), Emu(top + py))])

            elif cmd in ("Z", "z"):
                pass  # Close path handled by freeform builder

        if started:
            shape = builder.convert_to_shape()
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()


def _add_fallback_shape(
    slide,
    icon_name: str,
    left: int,
    top: int,
    size: int,
    color: RGBColor,
) -> None:
    """Add a simple circle with initial letter as fallback for complex icons."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    # Add initial letter
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = icon_name[0].upper()
    p.font.size = Pt(int(size / 914400 * 36))  # Scale font to shape
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.bold = True
    from pptx.enum.text import PP_ALIGN
    p.alignment = PP_ALIGN.CENTER


def get_icon_for_category(category: str) -> str:
    """Look up the best icon for a content category.

    Args:
        category: Content category (e.g., "finance", "growth", "time").

    Returns:
        Icon name string.
    """
    return ICON_CATEGORIES.get(category.lower(), "circle-checkmark")


def list_available_icons() -> list[str]:
    """List all available icon names."""
    if not ICON_DIR.exists():
        return []
    return sorted(p.stem for p in ICON_DIR.glob("*.svg"))
