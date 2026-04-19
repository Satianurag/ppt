"""Infographic renderers — python-pptx shape-based equivalents of PPT Master's
52 SVG chart templates (hugohe3/ppt-master templates/charts/).

Each renderer creates native editable shapes (MSO_SHAPE) instead of embedded
SVG images. This approach follows the research recommendation (Option C Hybrid):
- python-pptx for slide structure
- MSO_SHAPE for infographic elements
- add_chart() for data charts

Renderers implemented here:
  timeline, comparison_columns, pros_cons, cycle_diagram, hub_spoke,
  pyramid, numbered_steps, kpi_cards, funnel, vertical_list,
  matrix_2x2, icon_grid

Each function takes (slide, grid, items, accent_fn) and places shapes.
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from step4.font_sizing import auto_font_size_pt, auto_font_size_aggressive
from step4.theme import get_accent_color


def _add_shaped_textbox(
    slide, left: int, top: int, width: int, height: int,
    text: str, font_size: Pt, bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.CENTER,
    fill_color: RGBColor | None = None,
    font_color: RGBColor = RGBColor(0x33, 0x33, 0x33),
    shape_type: int = MSO_SHAPE.ROUNDED_RECTANGLE,
) -> None:
    """Add a shape with text — the building block for all infographics."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = font_color


def _accent_rgb(index: int) -> RGBColor:
    """Get a default accent color by index for infographics."""
    palette = [
        RGBColor(0x21, 0x96, 0xF3),  # Blue
        RGBColor(0x4C, 0xAF, 0x50),  # Green
        RGBColor(0xFF, 0x98, 0x00),  # Orange
        RGBColor(0xE9, 0x1E, 0x63),  # Pink
        RGBColor(0x9C, 0x27, 0xB0),  # Purple
        RGBColor(0x00, 0xBC, 0xD4),  # Cyan
        RGBColor(0xFF, 0x57, 0x22),  # Deep Orange
        RGBColor(0x60, 0x7D, 0x8B),  # Blue Grey
    ]
    return palette[index % len(palette)]


# ── Timeline ─────────────────────────────────────────────────────────
# Inspired by PPT Master: timeline.svg, roadmap_vertical.svg, chevron_process.svg

def render_timeline(slide, grid, items: list[dict], template_type=None) -> None:
    """Render a horizontal timeline with milestone markers.

    Items: [{"label": str, "description": str}, ...]
    """
    n = len(items)
    if n == 0:
        return

    usable_width = grid.content_width
    usable_left = grid.content_left
    y_center = grid.content_top + Inches(2.0)

    # Horizontal line
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        usable_left, y_center + Inches(0.15),
        usable_width, Inches(0.06),
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(0xBD, 0xBD, 0xBD)
    line_shape.line.fill.background()

    spacing = int(usable_width / max(n, 1))
    marker_size = Inches(0.35)

    for i, item in enumerate(items):
        cx = usable_left + int(spacing * (i + 0.5)) - int(marker_size / 2)

        # Circle marker
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, cx, y_center, marker_size, marker_size,
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = _accent_rgb(i)
        circle.line.fill.background()

        # Label above
        label = item.get("label", f"Step {i+1}")
        _add_shaped_textbox(
            slide, cx - Inches(0.5), y_center - Inches(0.9),
            Inches(1.5), Inches(0.5),
            label, Pt(12), bold=True,
            fill_color=None,
            font_color=RGBColor(0x33, 0x33, 0x33),
            shape_type=MSO_SHAPE.RECTANGLE,
        )

        # Description below
        desc = item.get("description", "")
        if desc:
            font_sz = auto_font_size_pt(desc, "tight_col")
            _add_shaped_textbox(
                slide, cx - Inches(0.5), y_center + Inches(0.6),
                Inches(1.5), Inches(1.2),
                desc, font_sz,
                fill_color=None,
                font_color=RGBColor(0x66, 0x66, 0x66),
                shape_type=MSO_SHAPE.RECTANGLE,
            )


# ── Comparison Columns ───────────────────────────────────────────────
# Inspired by PPT Master: comparison_columns.svg, comparison_table.svg

def render_comparison_columns(slide, grid, items: list[dict], template_type=None) -> None:
    """Render side-by-side comparison columns.

    Items: [{"title": str, "points": [str, ...]}, ...]
    """
    n = len(items)
    if n == 0:
        return

    col_gap = Inches(0.2)
    total_gap = int(col_gap * (n - 1))
    col_width = int((grid.content_width - total_gap) / n)
    col_height = grid.content_height

    for i, item in enumerate(items):
        left = grid.content_left + int(i * (col_width + col_gap))
        top = grid.content_top

        # Column header
        _add_shaped_textbox(
            slide, left, top, col_width, Inches(0.6),
            item.get("title", f"Option {i+1}"),
            Pt(16), bold=True,
            fill_color=_accent_rgb(i),
            font_color=RGBColor(0xFF, 0xFF, 0xFF),
        )

        # Column body with points
        points = item.get("points", [])
        body_text = "\n".join(f"• {p}" for p in points)
        font_sz = auto_font_size_pt(body_text, "two_col")

        body = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top + Inches(0.7), col_width, Inches(4.0),
        )
        body.fill.solid()
        body.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        body.line.fill.background()

        tf = body.text_frame
        tf.word_wrap = True
        for j, point in enumerate(points):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = font_sz
            p.space_after = Pt(6)


# ── Pros/Cons ────────────────────────────────────────────────────────
# Inspired by PPT Master: pros_cons_chart.svg

def render_pros_cons(slide, grid, pros: list[str], cons: list[str], template_type=None) -> None:
    """Render a two-column pros vs cons layout."""
    col_width = int(grid.content_width / 2) - Inches(0.1)
    top = grid.content_top

    # Pros column (green)
    _add_shaped_textbox(
        slide, grid.content_left, top, col_width, Inches(0.5),
        "PROS", Pt(18), bold=True,
        fill_color=RGBColor(0x4C, 0xAF, 0x50),
        font_color=RGBColor(0xFF, 0xFF, 0xFF),
    )
    for i, pro in enumerate(pros):
        _add_shaped_textbox(
            slide,
            grid.content_left, top + Inches(0.6 + i * 0.6),
            col_width, Inches(0.5),
            f"+ {pro}", Pt(13),
            fill_color=RGBColor(0xE8, 0xF5, 0xE9),
            font_color=RGBColor(0x2E, 0x7D, 0x32),
            shape_type=MSO_SHAPE.RECTANGLE,
        )

    # Cons column (red)
    right_left = grid.content_left + col_width + Inches(0.2)
    _add_shaped_textbox(
        slide, right_left, top, col_width, Inches(0.5),
        "CONS", Pt(18), bold=True,
        fill_color=RGBColor(0xF4, 0x43, 0x36),
        font_color=RGBColor(0xFF, 0xFF, 0xFF),
    )
    for i, con in enumerate(cons):
        _add_shaped_textbox(
            slide,
            right_left, top + Inches(0.6 + i * 0.6),
            col_width, Inches(0.5),
            f"- {con}", Pt(13),
            fill_color=RGBColor(0xFF, 0xEB, 0xEE),
            font_color=RGBColor(0xC6, 0x28, 0x28),
            shape_type=MSO_SHAPE.RECTANGLE,
        )


# ── Cycle Diagram ────────────────────────────────────────────────────
# Inspired by PPT Master: cycle_diagram.svg

def render_cycle_diagram(slide, grid, items: list[str], template_type=None) -> None:
    """Render a circular cycle diagram with connected nodes."""
    import math

    n = len(items)
    if n == 0:
        return

    cx = int(grid.content_left + grid.content_width / 2)
    cy = int(grid.content_top + grid.content_height / 2) - Inches(0.3)
    radius = Inches(2.2)
    node_size = Inches(1.4)

    for i, item in enumerate(items):
        angle = (2 * math.pi * i / n) - math.pi / 2
        x = cx + int(radius * math.cos(angle)) - int(node_size / 2)
        y = cy + int(radius * math.sin(angle)) - int(node_size / 2)

        font_sz = auto_font_size_pt(item, "tight_col")
        _add_shaped_textbox(
            slide, x, y, node_size, node_size,
            item, font_sz, bold=True,
            fill_color=_accent_rgb(i),
            font_color=RGBColor(0xFF, 0xFF, 0xFF),
            shape_type=MSO_SHAPE.OVAL,
        )


# ── Hub & Spoke ──────────────────────────────────────────────────────
# Inspired by PPT Master: hub_spoke.svg

def render_hub_spoke(slide, grid, center: str, spokes: list[str], template_type=None) -> None:
    """Render a hub-and-spoke diagram."""
    import math

    cx = int(grid.content_left + grid.content_width / 2)
    cy = int(grid.content_top + grid.content_height / 2) - Inches(0.3)

    # Center hub
    hub_size = Inches(1.8)
    _add_shaped_textbox(
        slide,
        cx - int(hub_size / 2), cy - int(hub_size / 2),
        hub_size, hub_size,
        center, Pt(16), bold=True,
        fill_color=RGBColor(0x21, 0x96, 0xF3),
        font_color=RGBColor(0xFF, 0xFF, 0xFF),
        shape_type=MSO_SHAPE.OVAL,
    )

    # Spokes
    n = len(spokes)
    radius = Inches(2.5)
    spoke_size = Inches(1.3)

    for i, spoke in enumerate(spokes):
        angle = (2 * math.pi * i / n) - math.pi / 2
        x = cx + int(radius * math.cos(angle)) - int(spoke_size / 2)
        y = cy + int(radius * math.sin(angle)) - int(spoke_size / 2)

        font_sz = auto_font_size_pt(spoke, "tight_col")
        _add_shaped_textbox(
            slide, x, y, spoke_size, spoke_size,
            spoke, font_sz,
            fill_color=_accent_rgb(i),
            font_color=RGBColor(0xFF, 0xFF, 0xFF),
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        )


# ── Pyramid ──────────────────────────────────────────────────────────
# Inspired by PPT Master: pyramid_chart.svg, isometric_stairs.svg

def render_pyramid(slide, grid, items: list[str], template_type=None) -> None:
    """Render a pyramid diagram (widest at bottom)."""
    n = len(items)
    if n == 0:
        return

    total_height = grid.content_height
    layer_height = int(total_height / n) - Inches(0.05)
    max_width = grid.content_width
    center_x = grid.content_left + int(grid.content_width / 2)

    for i, item in enumerate(items):
        fraction = (n - i) / n
        layer_width = int(max_width * fraction * 0.85)
        left = center_x - int(layer_width / 2)
        top = grid.content_top + int(i * (layer_height + Inches(0.05)))

        font_sz = auto_font_size_pt(item, "two_col")
        _add_shaped_textbox(
            slide, left, top, layer_width, layer_height,
            item, font_sz, bold=True,
            fill_color=_accent_rgb(i),
            font_color=RGBColor(0xFF, 0xFF, 0xFF),
            shape_type=MSO_SHAPE.RECTANGLE,
        )


# ── Numbered Steps ───────────────────────────────────────────────────
# Inspired by PPT Master: numbered_steps.svg, snake_flow.svg

def render_numbered_steps(slide, grid, items: list[dict], template_type=None) -> None:
    """Render numbered step cards in a row or grid.

    Items: [{"title": str, "description": str}, ...]
    """
    n = len(items)
    if n == 0:
        return

    # Dynamic layout: 1-3 in a row, 4+ in 2 rows
    if n <= 3:
        rows, cols = 1, n
    elif n <= 6:
        rows = 2
        cols = (n + 1) // 2
    else:
        rows = 2
        cols = (n + 1) // 2

    col_gap = Inches(0.2)
    row_gap = Inches(0.3)
    card_width = int((grid.content_width - col_gap * (cols - 1)) / cols)
    card_height = int((grid.content_height - row_gap * (rows - 1)) / rows)

    for idx, item in enumerate(items):
        row = idx // cols
        col = idx % cols

        left = grid.content_left + int(col * (card_width + col_gap))
        top = grid.content_top + int(row * (card_height + row_gap))

        # Number circle
        num_size = Inches(0.45)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left, top, num_size, num_size,
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = _accent_rgb(idx)
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = str(idx + 1)
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Title
        title = item.get("title", f"Step {idx+1}")
        _add_shaped_textbox(
            slide, left + Inches(0.5), top,
            card_width - Inches(0.5), Inches(0.4),
            title, Pt(14), bold=True,
            fill_color=None,
            font_color=RGBColor(0x33, 0x33, 0x33),
            shape_type=MSO_SHAPE.RECTANGLE,
        )

        # Description
        desc = item.get("description", "")
        if desc:
            font_sz = auto_font_size_pt(desc, "tight_col")
            _add_shaped_textbox(
                slide, left, top + Inches(0.5),
                card_width, card_height - Inches(0.6),
                desc, font_sz,
                fill_color=RGBColor(0xF5, 0xF5, 0xF5),
                font_color=RGBColor(0x55, 0x55, 0x55),
            )


# ── KPI Cards ────────────────────────────────────────────────────────
# Inspired by PPT Master: kpi_cards.svg, gauge_chart.svg

def render_kpi_cards(slide, grid, items: list[dict], template_type=None) -> None:
    """Render KPI/metric cards in a row.

    Items: [{"value": str, "label": str, "unit": str?}, ...]
    """
    n = len(items)
    if n == 0:
        return

    col_gap = Inches(0.25)
    card_width = int((grid.content_width - col_gap * (n - 1)) / n)
    card_height = Inches(3.0)
    top = grid.content_top + Inches(1.0)

    for i, item in enumerate(items):
        left = grid.content_left + int(i * (card_width + col_gap))

        # Card background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, card_width, card_height,
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

        # Accent bar at top
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, card_width, Inches(0.08),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = _accent_rgb(i)
        accent.line.fill.background()

        # Value (large)
        value = str(item.get("value", "—"))
        unit = item.get("unit", "")
        display_value = f"{value}{unit}" if unit else value
        _add_shaped_textbox(
            slide, left, top + Inches(0.5),
            card_width, Inches(1.2),
            display_value, Pt(36), bold=True,
            fill_color=None,
            font_color=_accent_rgb(i),
            shape_type=MSO_SHAPE.RECTANGLE,
        )

        # Label
        label = item.get("label", "")
        _add_shaped_textbox(
            slide, left, top + Inches(1.8),
            card_width, Inches(0.6),
            label, Pt(14),
            fill_color=None,
            font_color=RGBColor(0x66, 0x66, 0x66),
            shape_type=MSO_SHAPE.RECTANGLE,
        )


# ── Funnel ───────────────────────────────────────────────────────────
# Inspired by PPT Master: funnel_chart.svg

def render_funnel(slide, grid, items: list[dict], template_type=None) -> None:
    """Render a funnel diagram (widest at top, narrowest at bottom).

    Items: [{"label": str, "value": str?}, ...]
    """
    n = len(items)
    if n == 0:
        return

    max_width = grid.content_width
    layer_height = min(Inches(0.8), int(grid.content_height / n) - Inches(0.05))
    center_x = grid.content_left + int(grid.content_width / 2)

    for i, item in enumerate(items):
        fraction = 1.0 - (i * 0.12)
        layer_width = int(max_width * max(fraction, 0.3))
        left = center_x - int(layer_width / 2)
        top = grid.content_top + int(i * (layer_height + Inches(0.05)))

        label = item.get("label", "")
        value = item.get("value", "")
        text = f"{label}: {value}" if value else label

        font_sz = auto_font_size_pt(text, "two_col")
        _add_shaped_textbox(
            slide, left, top, layer_width, layer_height,
            text, font_sz, bold=True,
            fill_color=_accent_rgb(i),
            font_color=RGBColor(0xFF, 0xFF, 0xFF),
            shape_type=MSO_SHAPE.RECTANGLE,
        )


# ── Vertical List ────────────────────────────────────────────────────
# Inspired by PPT Master: vertical_list.svg

def render_vertical_list(slide, grid, items: list[dict], template_type=None) -> None:
    """Render a vertical list with accent markers.

    Items: [{"title": str, "description": str?}, ...]
    """
    n = len(items)
    if n == 0:
        return

    item_height = min(Inches(0.8), int(grid.content_height / n) - Inches(0.1))
    marker_width = Inches(0.08)

    for i, item in enumerate(items):
        top = grid.content_top + int(i * (item_height + Inches(0.1)))

        # Accent marker
        marker = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            grid.content_left, top, marker_width, item_height,
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = _accent_rgb(i)
        marker.line.fill.background()

        # Title
        title = item.get("title", "")
        _add_shaped_textbox(
            slide, grid.content_left + Inches(0.2), top,
            int(grid.content_width * 0.3), item_height,
            title, Pt(14), bold=True,
            fill_color=None, alignment=PP_ALIGN.LEFT,
            font_color=RGBColor(0x33, 0x33, 0x33),
            shape_type=MSO_SHAPE.RECTANGLE,
        )

        # Description
        desc = item.get("description", "")
        if desc:
            font_sz = auto_font_size_pt(desc, "body")
            _add_shaped_textbox(
                slide, grid.content_left + int(grid.content_width * 0.35), top,
                int(grid.content_width * 0.65), item_height,
                desc, font_sz,
                fill_color=None, alignment=PP_ALIGN.LEFT,
                font_color=RGBColor(0x55, 0x55, 0x55),
                shape_type=MSO_SHAPE.RECTANGLE,
            )


# ── Matrix 2×2 ───────────────────────────────────────────────────────
# Inspired by PPT Master: matrix_2x2.svg, swot_analysis.svg

def render_matrix_2x2(slide, grid, quadrants: list[dict], template_type=None) -> None:
    """Render a 2×2 matrix (SWOT, BCG, risk matrix, etc.).

    Quadrants: [{"title": str, "points": [str, ...]}, ...] (exactly 4)
    """
    if len(quadrants) < 4:
        quadrants = quadrants + [{"title": "", "points": []}] * (4 - len(quadrants))

    gap = Inches(0.15)
    cell_width = int((grid.content_width - gap) / 2)
    cell_height = int((grid.content_height - gap) / 2)

    positions = [
        (grid.content_left, grid.content_top),
        (grid.content_left + cell_width + gap, grid.content_top),
        (grid.content_left, grid.content_top + cell_height + gap),
        (grid.content_left + cell_width + gap, grid.content_top + cell_height + gap),
    ]

    for i, (left, top) in enumerate(positions):
        q = quadrants[i]

        # Header
        _add_shaped_textbox(
            slide, left, top, cell_width, Inches(0.5),
            q.get("title", ""), Pt(14), bold=True,
            fill_color=_accent_rgb(i),
            font_color=RGBColor(0xFF, 0xFF, 0xFF),
        )

        # Points
        points = q.get("points", [])
        body_text = "\n".join(f"• {p}" for p in points)
        font_sz = auto_font_size_pt(body_text, "tight_col")

        body = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top + Inches(0.55), cell_width, cell_height - Inches(0.6),
        )
        body.fill.solid()
        body.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        body.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

        tf = body.text_frame
        tf.word_wrap = True
        for j, point in enumerate(points):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = font_sz
            p.space_after = Pt(4)


# ── Icon Grid ────────────────────────────────────────────────────────
# Inspired by PPT Master: icon_grid.svg

def render_icon_grid(slide, grid, items: list[dict], template_type=None) -> None:
    """Render an icon-style grid (using emoji/text as icon substitute).

    Items: [{"icon": str, "title": str, "description": str?}, ...]
    """
    n = len(items)
    if n == 0:
        return

    if n <= 3:
        rows, cols = 1, n
    elif n <= 6:
        rows, cols = 2, min(n, 3)
    else:
        rows, cols = 3, min(n, 3)

    col_gap = Inches(0.3)
    row_gap = Inches(0.3)
    card_width = int((grid.content_width - col_gap * (cols - 1)) / cols)
    card_height = int((grid.content_height - row_gap * (rows - 1)) / rows)

    for idx, item in enumerate(items):
        row = idx // cols
        col = idx % cols
        if row >= rows:
            break

        left = grid.content_left + int(col * (card_width + col_gap))
        top = grid.content_top + int(row * (card_height + row_gap))

        # Icon (text substitute — large character)
        icon = item.get("icon", "●")
        _add_shaped_textbox(
            slide, left, top, card_width, Inches(0.6),
            icon, Pt(28), bold=True,
            fill_color=None,
            font_color=_accent_rgb(idx),
            shape_type=MSO_SHAPE.RECTANGLE,
        )

        # Title
        title = item.get("title", "")
        _add_shaped_textbox(
            slide, left, top + Inches(0.6),
            card_width, Inches(0.4),
            title, Pt(14), bold=True,
            fill_color=None,
            font_color=RGBColor(0x33, 0x33, 0x33),
            shape_type=MSO_SHAPE.RECTANGLE,
        )

        # Description
        desc = item.get("description", "")
        if desc:
            font_sz = auto_font_size_pt(desc, "tight_col")
            _add_shaped_textbox(
                slide, left, top + Inches(1.0),
                card_width, card_height - Inches(1.1),
                desc, font_sz,
                fill_color=None, alignment=PP_ALIGN.CENTER,
                font_color=RGBColor(0x66, 0x66, 0x66),
                shape_type=MSO_SHAPE.RECTANGLE,
            )


# ── Dynamic Layout Selection ─────────────────────────────────────────

INFOGRAPHIC_RENDERERS = {
    "timeline": render_timeline,
    "comparison_columns": render_comparison_columns,
    "pros_cons": render_pros_cons,
    "cycle_diagram": render_cycle_diagram,
    "hub_spoke": render_hub_spoke,
    "pyramid": render_pyramid,
    "numbered_steps": render_numbered_steps,
    "kpi_cards": render_kpi_cards,
    "funnel": render_funnel,
    "vertical_list": render_vertical_list,
    "matrix_2x2": render_matrix_2x2,
    "icon_grid": render_icon_grid,
}
