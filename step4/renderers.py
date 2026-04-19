"""Individual slide renderers: cover, end, bullets, charts, tables, infographics.

Reusable patterns:
- PPTAgent apis.py:329-341 — table auto-column-width
- PPTAgent apis.py:344-352 — table cell merging
- solution_architecture.md GAP 3 — table styling (zebra stripes, bold headers)
- hackathon_final_analysis1.md §9.1 — chart theme color assignment
- gap_analysis1.md VIZ-2 — single data point chart → stat card
- gap_analysis1.md VIZ-7 — chevron shapes for process flows

No images rendered (user explicitly excluded images from scope).
"""

import re
from datetime import datetime

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from step2.slide_plan_models import SlideType, LayoutType, ChartType
from step3.content_models import SlideContent, ChartData, TableData
from step4.grid import Grid, TITLE_FONT_SIZE, SUBTITLE_FONT_SIZE, BODY_FONT_SIZE, FOOTER_FONT_SIZE
from step4.template_manager import (
    TemplateType, LayoutRole, get_layout, COVER_PLACEHOLDERS,
)
from step4.theme import get_accent_color


# python-pptx chart type mapping
CHART_TYPE_MAP: dict[ChartType, int] = {
    ChartType.BAR: XL_CHART_TYPE.COLUMN_CLUSTERED,
    ChartType.HORIZONTAL_BAR: XL_CHART_TYPE.BAR_CLUSTERED,
    ChartType.GROUPED_BAR: XL_CHART_TYPE.COLUMN_CLUSTERED,
    ChartType.LINE: XL_CHART_TYPE.LINE_MARKERS,
    ChartType.PIE: XL_CHART_TYPE.PIE,
    ChartType.DONUT: XL_CHART_TYPE.DOUGHNUT,
}

# Minimum data points for chart types (VIZ-2)
MIN_DATA_POINTS: dict[ChartType, int] = {
    ChartType.BAR: 2,
    ChartType.HORIZONTAL_BAR: 2,
    ChartType.GROUPED_BAR: 2,
    ChartType.LINE: 2,
    ChartType.PIE: 2,
    ChartType.DONUT: 2,
}


def _add_textbox(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    font_size: Pt = BODY_FONT_SIZE,
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    color: RGBColor | None = None,
) -> None:
    """Add a text box shape to a slide. Never sets font.name (GAP-11: let theme cascade)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.alignment = alignment
    if color is not None:
        p.font.color.rgb = color


def _add_title_and_subtitle(
    slide, grid: Grid, title: str, subtitle: str | None = None
) -> None:
    """Add title and optional subtitle text boxes to a content slide."""
    _add_textbox(
        slide,
        grid.title_left, grid.title_top,
        grid.title_width, grid.title_height,
        title,
        font_size=TITLE_FONT_SIZE,
        bold=True,
    )
    if subtitle:
        _add_textbox(
            slide,
            grid.title_left, grid.subtitle_top,
            grid.title_width, grid.subtitle_height,
            subtitle,
            font_size=SUBTITLE_FONT_SIZE,
            color=RGBColor(0x66, 0x66, 0x66),
        )


def _set_uae_title_placeholder(slide, title: str) -> None:
    """Set the title placeholder (idx=0) on UAE Solar layout slides."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            break


def _add_slide_title(slide, grid: Grid, template_type: TemplateType, title: str, subtitle: str | None = None) -> None:
    """Add title to a slide — uses placeholder for UAE Solar, textbox for others."""
    if template_type == TemplateType.UAE_SOLAR:
        _set_uae_title_placeholder(slide, title)
    else:
        _add_title_and_subtitle(slide, grid, title, subtitle)


def _add_key_message_footer(slide, grid: Grid, key_message: str) -> None:
    """Add a key message footer at the bottom of a content slide."""
    if not key_message:
        return
    _add_textbox(
        slide,
        grid.footer_left, grid.footer_top,
        grid.footer_width, grid.footer_height,
        key_message,
        font_size=FOOTER_FONT_SIZE,
        color=RGBColor(0x99, 0x99, 0x99),
        alignment=PP_ALIGN.LEFT,
    )


def _detect_date_from_title(title: str) -> str:
    """Try to extract a date from the presentation title, else use today."""
    date_pattern = re.compile(r'(\d{4}[-/]\d{2}[-/]\d{2}|\d{8})')
    match = date_pattern.search(title)
    if match:
        raw = match.group(1).replace("/", "-")
        if len(raw) == 8 and raw.isdigit():
            return f"{raw[:4]}-{raw[4:6]}-{raw[6:]}"
        return raw
    return datetime.now().strftime("%B %d, %Y")


# ── Cover slide ──────────────────────────────────────────────────────

def render_cover(
    prs: Presentation,
    template_type: TemplateType,
    slide_content: SlideContent,
    presentation_title: str,
) -> None:
    """Render the cover/title slide using the template's existing cover layout.

    User's idea: use the template's designed cover as-is, just populate placeholders
    with title, subtitle, presenter name, and date.
    """
    layout = get_layout(prs, template_type, LayoutRole.COVER)
    slide = prs.slides.add_slide(layout)

    ph_map = COVER_PLACEHOLDERS.get(template_type, {"title": 10, "subtitle": 11})

    # Populate title placeholder
    title_idx = ph_map["title"]
    if title_idx is not None:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == title_idx:
                ph.text = presentation_title
                for p in ph.text_frame.paragraphs:
                    p.font.size = TITLE_FONT_SIZE
                break

    # Populate subtitle placeholder (subtitle + date)
    subtitle_idx = ph_map.get("subtitle")
    date_str = _detect_date_from_title(presentation_title)
    subtitle_text = slide_content.subtitle or ""
    if subtitle_text:
        subtitle_text = f"{subtitle_text}\n{date_str}"
    else:
        subtitle_text = date_str

    if subtitle_idx is not None:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == subtitle_idx:
                ph.text = subtitle_text
                for p in ph.text_frame.paragraphs:
                    p.font.size = SUBTITLE_FONT_SIZE
                break
    else:
        # UAE Solar: no subtitle placeholder → add a text box manually
        _add_textbox(
            slide,
            Inches(1.0), Inches(4.5),
            Inches(11.33), Inches(0.5),
            subtitle_text,
            font_size=SUBTITLE_FONT_SIZE,
            alignment=PP_ALIGN.CENTER,
            color=RGBColor(0x66, 0x66, 0x66),
        )


# ── Thank-you / End slide ───────────────────────────────────────────

def render_end(
    prs: Presentation,
    template_type: TemplateType,
    slide_content: SlideContent,
) -> None:
    """Render the closing slide using the template's existing thank-you layout.

    User's idea: the template already says "Thank You" — just add it as-is.
    """
    layout = get_layout(prs, template_type, LayoutRole.END)
    prs.slides.add_slide(layout)


# ── Bullet slide ─────────────────────────────────────────────────────

def render_bullets(
    prs: Presentation,
    template_type: TemplateType,
    slide_content: SlideContent,
    grid: Grid,
) -> None:
    """Render a bullet-point content slide on the Blank/content layout."""
    layout = get_layout(prs, template_type, LayoutRole.CONTENT)
    slide = prs.slides.add_slide(layout)

    # For UAE Solar (layout 2 has placeholders), use placeholders for title/subtitle
    # but render bullets in a textbox (idx=11 is a tiny unusable footnote area)
    if template_type == TemplateType.UAE_SOLAR:
        _render_bullets_uae(slide, slide_content, grid)
        _add_key_message_footer(slide, grid, slide_content.key_message)
        return

    _add_title_and_subtitle(slide, grid, slide_content.title, slide_content.subtitle)

    # Build bullet text box
    txBox = slide.shapes.add_textbox(
        grid.content_left, grid.content_top,
        grid.content_width, grid.content_height,
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    bullets = _get_bullet_texts(slide_content)
    for i, bullet_text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet_text}"
        p.font.size = BODY_FONT_SIZE
        p.space_after = Pt(8)

    _add_key_message_footer(slide, grid, slide_content.key_message)


def _render_bullets_uae(slide, slide_content: SlideContent, grid: Grid) -> None:
    """UAE Solar uses Layout 2 placeholders for title/subtitle only.

    placeholder idx=11 is a 0.13"-tall footnote at y=6.74" — unusable for
    bullet content. Instead, populate title/subtitle placeholders and render
    bullets in a programmatic textbox in the safe content area.
    """
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:  # Title
            ph.text = slide_content.title
        elif idx == 2:  # Subtitle
            ph.text = slide_content.subtitle or ""

    # Render bullets in a textbox (not in the tiny idx=11 footnote placeholder)
    bullets = _get_bullet_texts(slide_content)
    if bullets:
        txBox = slide.shapes.add_textbox(
            grid.content_left, grid.content_top,
            grid.content_width, grid.content_height,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet_text in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {bullet_text}"
            p.font.size = BODY_FONT_SIZE
            p.space_after = Pt(8)


def _get_bullet_texts(slide_content: SlideContent) -> list[str]:
    """Extract bullet text from either bullets or key_points."""
    if slide_content.bullets:
        return [b.text for b in slide_content.bullets]
    if slide_content.key_points:
        texts = []
        for kp in slide_content.key_points:
            texts.extend(kp.bullet_form)
        return texts[:6]  # MAX_BULLETS_PER_SLIDE
    return []


# ── Chart slide ──────────────────────────────────────────────────────

def render_chart(
    prs: Presentation,
    template_type: TemplateType,
    slide_content: SlideContent,
    grid: Grid,
) -> None:
    """Render a native python-pptx chart on a content slide.

    Charts do NOT auto-inherit theme colors — must manually set
    MSO_THEME_COLOR.ACCENT_N per series.
    """
    chart_data_model = slide_content.chart_data
    if chart_data_model is None:
        render_bullets(prs, template_type, slide_content, grid)
        return

    # VIZ-2: Check minimum data points — render stat card if too few
    if len(chart_data_model.categories) < MIN_DATA_POINTS.get(chart_data_model.chart_type, 2):
        _render_stat_card(prs, template_type, slide_content, grid, chart_data_model)
        return

    layout = get_layout(prs, template_type, LayoutRole.CONTENT)
    slide = prs.slides.add_slide(layout)

    _add_slide_title(slide, grid, template_type, slide_content.title)

    # Chart creation — shared across all templates
    chart_pptx_data = CategoryChartData()
    chart_pptx_data.categories = chart_data_model.categories
    for series_dict in chart_data_model.series:
        chart_pptx_data.add_series(series_dict["name"], series_dict["values"])

    xl_type = CHART_TYPE_MAP.get(chart_data_model.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    chart_shape = slide.shapes.add_chart(
        xl_type,
        grid.chart_left, grid.chart_top,
        grid.chart_width, grid.chart_height,
        chart_pptx_data,
    )
    _style_chart(chart_shape.chart, chart_data_model)
    _add_key_message_footer(slide, grid, slide_content.key_message)


def _style_chart(chart, chart_data_model: ChartData) -> None:
    """Apply theme colors and formatting to a chart.

    Pie/donut charts have ONE series with multiple data points — must color
    individual points, not the series. Bar/line charts color per series.
    """
    chart.has_legend = chart_data_model.show_legend

    is_single_series_chart = chart_data_model.chart_type in (ChartType.PIE, ChartType.DONUT)

    if is_single_series_chart:
        # Pie/donut: color each data point (slice) individually
        series = chart.series[0]
        for i in range(len(chart_data_model.categories)):
            point = series.points[i]
            point.format.fill.solid()
            point.format.fill.fore_color.theme_color = get_accent_color(i)
    else:
        # Bar/line/grouped: color each series
        for i, series in enumerate(chart.series):
            series.format.fill.solid()
            series.format.fill.fore_color.theme_color = get_accent_color(i)

    if chart_data_model.show_data_labels:
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(9)
        data_labels.number_format = chart_data_model.number_format


def _render_stat_card(
    prs: Presentation,
    template_type: TemplateType,
    slide_content: SlideContent,
    grid: Grid,
    chart_data_model: ChartData,
) -> None:
    """Render a stat card for charts with too few data points (VIZ-2)."""
    layout = get_layout(prs, template_type, LayoutRole.CONTENT)
    slide = prs.slides.add_slide(layout)

    _add_slide_title(slide, grid, template_type, slide_content.title)

    # Render each data point as a large number with label
    if chart_data_model.series and chart_data_model.categories:
        n_items = len(chart_data_model.categories)
        card_width = Inches(min(3.0, 10.0 / max(n_items, 1)))
        card_height = Inches(2.5)

        for i, category in enumerate(chart_data_model.categories):
            left = grid.content_left + Inches(i * 3.5)
            top = Inches(2.5)

            value = chart_data_model.series[0]["values"][i] if chart_data_model.series[0]["values"] else 0

            # Large number
            _add_textbox(
                slide, left, top, card_width, Inches(1.2),
                str(value), font_size=Pt(36), bold=True,
                alignment=PP_ALIGN.CENTER,
            )
            # Label below
            _add_textbox(
                slide, left, top + Inches(1.3), card_width, Inches(0.5),
                category, font_size=Pt(14),
                alignment=PP_ALIGN.CENTER,
                color=RGBColor(0x66, 0x66, 0x66),
            )

    _add_key_message_footer(slide, grid, slide_content.key_message)


# ── Table slide ──────────────────────────────────────────────────────

def render_table(
    prs: Presentation,
    template_type: TemplateType,
    slide_content: SlideContent,
    grid: Grid,
) -> None:
    """Render a styled table on a content slide.

    Reuses PPTAgent's add_table() auto-column-width pattern (apis.py:329-341)
    and adds custom zebra stripes + bold header styling (solution_architecture.md GAP 3).
    """
    table_data = slide_content.table_data
    if table_data is None:
        render_bullets(prs, template_type, slide_content, grid)
        return

    layout = get_layout(prs, template_type, LayoutRole.CONTENT)
    slide = prs.slides.add_slide(layout)

    _add_slide_title(slide, grid, template_type, slide_content.title)

    rows = len(table_data.rows) + 1  # +1 for header
    cols = len(table_data.headers)

    # Limit table size for readability
    max_rows = min(rows, 12)
    max_cols = min(cols, 8)

    table_shape = slide.shapes.add_table(
        max_rows, max_cols,
        grid.table_left, grid.table_top,
        grid.table_width, grid.table_height,
    )
    table = table_shape.table

    # Auto column width based on content length (PPTAgent pattern)
    all_data = [table_data.headers[:max_cols]] + [r[:max_cols] for r in table_data.rows[:max_rows - 1]]
    max_lengths = [
        max(len(str(all_data[r][c])) for r in range(len(all_data)))
        for c in range(max_cols)
    ]
    total_length = max(sum(max_lengths), 1)
    for c in range(max_cols):
        table.columns[c].width = int((max_lengths[c] / total_length) * grid.table_width)

    # Fill header row
    for c in range(max_cols):
        table.cell(0, c).text = str(table_data.headers[c])

    # Fill data rows
    for r in range(min(len(table_data.rows), max_rows - 1)):
        for c in range(max_cols):
            cell_val = table_data.rows[r][c] if c < len(table_data.rows[r]) else ""
            table.cell(r + 1, c).text = str(cell_val)

    # Apply styling
    _style_table(table, max_rows, max_cols, table_data)

    _add_key_message_footer(slide, grid, slide_content.key_message)


def _style_table(table, rows: int, cols: int, table_data: TableData) -> None:
    """Apply zebra stripes, bold headers, theme colors (solution_architecture.md GAP 3)."""
    font_size = Pt(table_data.recommended_font_size)
    zebra_light = RGBColor(0xF0, 0xF0, 0xF0)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = font_size
                paragraph.alignment = PP_ALIGN.CENTER

                # Right-align numeric columns
                if c in table_data.has_numeric_columns:
                    paragraph.alignment = PP_ALIGN.RIGHT

            if r == 0 and table_data.bold_headers:
                cell.fill.solid()
                cell.fill.fore_color.theme_color = get_accent_color(0)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            elif r % 2 == 0 and table_data.zebra_stripes:
                cell.fill.solid()
                cell.fill.fore_color.rgb = zebra_light

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


# ── Infographic slide (process flow, timeline, comparison) ───────────

def render_infographic(
    prs: Presentation,
    template_type: TemplateType,
    slide_content: SlideContent,
    grid: Grid,
) -> None:
    """Render infographic slides using MSO_SHAPE shapes.

    Uses chevrons for process flows (VIZ-7: renders consistently across viewers).
    """
    if slide_content.layout == LayoutType.COMPARISON:
        _render_comparison(prs, template_type, slide_content, grid)
    elif slide_content.layout == LayoutType.TIMELINE:
        _render_timeline(prs, template_type, slide_content, grid)
    elif slide_content.layout == LayoutType.PROCESS:
        _render_process_flow(prs, template_type, slide_content, grid)
    else:
        _render_process_flow(prs, template_type, slide_content, grid)


def _render_process_flow(
    prs: Presentation,
    template_type: TemplateType,
    slide_content: SlideContent,
    grid: Grid,
) -> None:
    """Horizontal chevron process flow."""
    layout = get_layout(prs, template_type, LayoutRole.CONTENT)
    slide = prs.slides.add_slide(layout)

    _add_slide_title(slide, grid, template_type, slide_content.title)

    bullets = _get_bullet_texts(slide_content)
    n = min(len(bullets), 5)  # Max 5 steps in process flow
    if n == 0:
        return

    chevron_width = Inches(min(2.2, 11.0 / n))
    chevron_height = Inches(1.5)
    gap = Inches(0.15)
    start_left = grid.content_left
    top = grid.process_top

    for i in range(n):
        left = start_left + int(i * (chevron_width + gap))

        shape = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            left, top, chevron_width, chevron_height,
        )
        shape.fill.solid()
        shape.fill.fore_color.theme_color = get_accent_color(i)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].text = bullets[i]
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf.paragraphs[0].font.bold = True

    _add_key_message_footer(slide, grid, slide_content.key_message)


def _render_comparison(
    prs: Presentation,
    template_type: TemplateType,
    slide_content: SlideContent,
    grid: Grid,
) -> None:
    """Two-column comparison layout."""
    layout = get_layout(prs, template_type, LayoutRole.CONTENT)
    slide = prs.slides.add_slide(layout)

    _add_slide_title(slide, grid, template_type, slide_content.title)

    bullets = _get_bullet_texts(slide_content)
    mid = len(bullets) // 2

    # Left column
    left_box = slide.shapes.add_textbox(
        grid.left_col_left, grid.content_top,
        grid.left_col_width, grid.content_height,
    )
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    for i, text in enumerate(bullets[:mid]):
        p = tf_left.paragraphs[0] if i == 0 else tf_left.add_paragraph()
        p.text = f"• {text}"
        p.font.size = BODY_FONT_SIZE
        p.space_after = Pt(6)

    # Right column
    right_box = slide.shapes.add_textbox(
        grid.right_col_left, grid.content_top,
        grid.right_col_width, grid.content_height,
    )
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    for i, text in enumerate(bullets[mid:]):
        p = tf_right.paragraphs[0] if i == 0 else tf_right.add_paragraph()
        p.text = f"• {text}"
        p.font.size = BODY_FONT_SIZE
        p.space_after = Pt(6)

    _add_key_message_footer(slide, grid, slide_content.key_message)


def _render_timeline(
    prs: Presentation,
    template_type: TemplateType,
    slide_content: SlideContent,
    grid: Grid,
) -> None:
    """Horizontal timeline with rounded rectangles."""
    layout = get_layout(prs, template_type, LayoutRole.CONTENT)
    slide = prs.slides.add_slide(layout)

    _add_slide_title(slide, grid, template_type, slide_content.title)

    bullets = _get_bullet_texts(slide_content)
    n = min(len(bullets), 6)
    if n == 0:
        return

    node_width = Inches(min(1.8, 10.0 / n))
    node_height = Inches(1.0)
    line_y = Inches(3.0)  # Center line of timeline
    gap = Inches(0.2)

    # Draw horizontal line
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        grid.content_left, line_y + Inches(0.4),
        grid.content_width, Inches(0.05),
    ).fill.solid()

    for i in range(n):
        left = grid.content_left + int(i * (node_width + gap))
        # Alternate above/below the line
        top = line_y - Inches(0.8) if i % 2 == 0 else line_y + Inches(0.8)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, node_width, node_height,
        )
        shape.fill.solid()
        shape.fill.fore_color.theme_color = get_accent_color(i)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].text = bullets[i]
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    _add_key_message_footer(slide, grid, slide_content.key_message)


# ── Agenda slide ─────────────────────────────────────────────────────

def render_agenda(
    prs: Presentation,
    template_type: TemplateType,
    slide_content: SlideContent,
    grid: Grid,
    all_slides: list[SlideContent],
) -> None:
    """Render an agenda/table of contents slide listing all section titles."""
    layout = get_layout(prs, template_type, LayoutRole.CONTENT)
    slide = prs.slides.add_slide(layout)

    _add_slide_title(slide, grid, template_type, "Agenda")

    if template_type != TemplateType.UAE_SOLAR:
        txBox = slide.shapes.add_textbox(
            grid.content_left, grid.content_top,
            grid.content_width, grid.content_height,
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        section_titles = [
            s.title for s in all_slides
            if s.slide_type not in (SlideType.TITLE, SlideType.AGENDA, SlideType.THANK_YOU)
        ]
        for i, title in enumerate(section_titles):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{i + 1}.  {title}"
            p.font.size = Pt(18)
            p.space_after = Pt(6)
    else:
        # UAE Solar: title already set by _add_slide_title; render agenda items in textbox
        section_titles = [
            s.title for s in all_slides
            if s.slide_type not in (SlideType.TITLE, SlideType.AGENDA, SlideType.THANK_YOU)
        ]
        if section_titles:
            txBox = slide.shapes.add_textbox(
                grid.content_left, grid.content_top,
                grid.content_width, grid.content_height,
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, title in enumerate(section_titles):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"{i + 1}.  {title}"
                p.font.size = Pt(18)
                p.space_after = Pt(6)


