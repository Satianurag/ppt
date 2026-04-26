"""Audit Canva/master PPTX files for editable, picture-free eligibility."""

from __future__ import annotations

import json
import re
import zipfile
from dataclasses import asdict, dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


EMU_PER_INCH = 914400


@dataclass
class SlideAudit:
    file: str
    slide: int
    width_in: float
    height_in: float
    top_level_shapes: int
    text_shapes: int
    picture_shapes: int
    chart_shapes: int
    media_relationships: int
    svg_media: int
    raster_media: int
    eligible: bool
    reason: str


def audit_pptx(path: str | Path) -> list[SlideAudit]:
    pptx = Path(path)
    prs = Presentation(str(pptx))
    rel_media = _slide_media_relationship_counts(pptx)
    media_names = _slide_media_names(pptx)
    rows: list[SlideAudit] = []
    for i, slide in enumerate(prs.slides, start=1):
        text_shapes = 0
        picture_shapes = 0
        chart_shapes = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_shapes += 1
            if _is_picture_shape(shape):
                picture_shapes += 1
            if shape.has_chart:
                chart_shapes += 1
        slide_media = media_names.get(i, [])
        svg_media = sum(1 for name in slide_media if name.lower().endswith(".svg"))
        raster_media = sum(
            1 for name in slide_media
            if name.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff"))
        )
        media_count = rel_media.get(i, 0)
        eligible = media_count == 0 and picture_shapes == 0
        reason = "native editable" if eligible else "contains picture/media"
        rows.append(
            SlideAudit(
                file=str(pptx),
                slide=i,
                width_in=round(prs.slide_width / EMU_PER_INCH, 3),
                height_in=round(prs.slide_height / EMU_PER_INCH, 3),
                top_level_shapes=len(slide.shapes),
                text_shapes=text_shapes,
                picture_shapes=picture_shapes,
                chart_shapes=chart_shapes,
                media_relationships=media_count,
                svg_media=svg_media,
                raster_media=raster_media,
                eligible=eligible,
                reason=reason,
            )
        )
    return rows


def audit_paths(paths: list[str | Path]) -> list[SlideAudit]:
    rows: list[SlideAudit] = []
    for path in paths:
        rows.extend(audit_pptx(path))
    return rows


def write_catalog(paths: list[str | Path], output_path: str | Path) -> str:
    rows = audit_paths(paths)
    data = {
        "slides": [asdict(row) for row in rows],
        "summary": {
            "files": len({row.file for row in rows}),
            "slides": len(rows),
            "eligible_picture_free": sum(1 for row in rows if row.eligible),
            "contains_picture_or_media": sum(1 for row in rows if not row.eligible),
        },
    }
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(data, indent=2), encoding="utf-8")
    return str(out)


def _is_picture_shape(shape) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.PICTURE


def _slide_media_relationship_counts(pptx: Path) -> dict[int, int]:
    counts: dict[int, int] = {}
    with zipfile.ZipFile(pptx) as zf:
        for name in zf.namelist():
            if not name.startswith("ppt/slides/_rels/slide") or not name.endswith(".xml.rels"):
                continue
            match = re.search(r"/slide(\d+)\.xml\.rels$", name)
            if match is None:
                continue
            slide_no = int(match.group(1))
            xml = zf.read(name).decode("utf-8", errors="ignore")
            counts[slide_no] = xml.count("/media/")
    return counts


def _slide_media_names(pptx: Path) -> dict[int, list[str]]:
    names: dict[int, list[str]] = {}
    with zipfile.ZipFile(pptx) as zf:
        for name in zf.namelist():
            if not name.startswith("ppt/slides/_rels/slide") or not name.endswith(".xml.rels"):
                continue
            match = re.search(r"/slide(\d+)\.xml\.rels$", name)
            if match is None:
                continue
            slide_no = int(match.group(1))
            xml = zf.read(name).decode("utf-8", errors="ignore")
            parts: list[str] = []
            for marker in ("../media/", "/media/"):
                fragments = xml.split(marker)[1:]
                for fragment in fragments:
                    parts.append(fragment.split('"', 1)[0])
            names[slide_no] = parts
    return names
