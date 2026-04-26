"""Post-render constraint checks.

Each constraint maps to a deterministic assertion executed against the rendered
.pptx file. The reviewer uses the returned issues to drive retry feedback.
"""

from __future__ import annotations

import json
import os
import re
import zipfile
from dataclasses import dataclass

from pptx import Presentation


_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


@dataclass
class AssertionResult:
    constraint: str
    passed: bool
    issues: list[str]


def check_c1_fixed_slides(pptx_path: str, expected_count: int = 15) -> AssertionResult:
    """C1: Exactly ``expected_count`` slides; cover + end preserve template layouts.

    The manifest records which layout class each slide uses; we rely on that
    instead of name-matching because UAE's thank-you layout shares a name with
    its body layouts.
    """
    prs = Presentation(pptx_path)
    issues: list[str] = []
    n = len(prs.slides)
    if n != expected_count:
        issues.append(f"C1: slide count {n} != {expected_count}")

    manifest = _load_manifest(pptx_path)
    if manifest:
        if not manifest or manifest[0]["class"] != "cover":
            issues.append("C1: slide 1 is not the template cover")
        if manifest and manifest[-1]["class"] != "end":
            issues.append("C1: last slide is not the template thank-you")
    else:
        if n >= 1:
            cover = prs.slides[0]
            if ("cover" not in cover.slide_layout.name.lower()
                    and "title" not in cover.slide_layout.name.lower()):
                issues.append(f"C1: slide 1 layout {cover.slide_layout.name!r} is not a cover layout")

    return AssertionResult("C1", not issues, issues)


def _layout_class_of(slide) -> str:
    """Classify a slide by its visible geometric shapes so C2 can compare layouts."""
    shape_bag: list[str] = []
    has_chart = False
    has_table = False
    rect_count = 0
    oval_count = 0
    chevron_count = 0
    trapezoid_count = 0
    textbox_count = 0

    for shape in slide.shapes:
        if shape.has_chart:
            has_chart = True
        if shape.has_table:
            has_table = True
        if shape.shape_type is None:
            continue
        name = str(shape.shape_type)
        if "CHART" in name:
            has_chart = True
        elif "TABLE" in name:
            has_table = True
        elif "OVAL" in name:
            oval_count += 1
        elif "CHEVRON" in name:
            chevron_count += 1
        elif "TRAPEZOID" in name or "TRIANGLE" in name:
            trapezoid_count += 1
        elif "ROUNDED_RECTANGLE" in name or "RECTANGLE" in name:
            rect_count += 1
        elif "TEXT_BOX" in name:
            textbox_count += 1

    if has_chart and rect_count >= 1:
        return "chart+text"
    if has_chart:
        return "chart"
    if has_table:
        return "table"
    if chevron_count >= 3:
        return "process"
    if trapezoid_count >= 2:
        return "pyramid"
    if oval_count >= 3:
        return "timeline"
    if rect_count >= 4:
        return "grid-4"
    if rect_count >= 3:
        return "grid-3"
    if rect_count == 2:
        return "compare"
    if rect_count == 1 and oval_count == 0:
        return "kpi"
    if textbox_count and not shape_bag:
        return "bullets"
    return "bullets"


def _load_manifest(pptx_path: str) -> list[dict] | None:
    manifest = pptx_path + ".layouts.json"
    if not os.path.exists(manifest):
        return None
    with open(manifest) as f:
        return json.load(f).get("layouts")


def check_c2_no_adjacent_repeat(pptx_path: str) -> AssertionResult:
    """C2: Body slides 2..14 must not share a geometric class with a neighbour.

    Uses the sidecar manifest produced by ``deck_builder`` — it records the
    exact layout class chosen for each slide, which is more reliable than
    reverse-engineering shape counts from the rendered PPTX.
    """
    issues: list[str] = []
    manifest = _load_manifest(pptx_path)
    if manifest:
        classes = [m["class"] for m in manifest]
    else:
        prs = Presentation(pptx_path)
        classes = [_layout_class_of(s) for s in prs.slides]

    # Check body range only (exclude cover + end)
    for i in range(1, len(classes) - 2):
        if classes[i] == classes[i + 1]:
            issues.append(
                f"C2: slides {i + 1} and {i + 2} share layout class {classes[i]!r}"
            )
    return AssertionResult("C2", not issues, issues)


def check_c3_master_inheritance(pptx_path: str) -> AssertionResult:
    """C3: No hardcoded srgbClr or a:latin font overrides inside slide XML."""
    issues: list[str] = []
    with zipfile.ZipFile(pptx_path) as zf:
        slide_names = [n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        for name in slide_names:
            xml = zf.read(name).decode("utf-8")
            # Count <a:srgbClr nodes *inside* a:solidFill or a:fill; these are
            # the ones that bypass master inheritance.
            # Charts are allowed to declare their own palette, so we exclude
            # drawings embedded as charts. Text-run srgbClr overrides are
            # forbidden though.
            bad_color = len(re.findall(r"<a:rPr[^>]*>.*?<a:solidFill><a:srgbClr", xml, re.DOTALL))
            if bad_color:
                issues.append(f"C3: {name} has {bad_color} hardcoded text colors")
            latin = len(re.findall(r"<a:rPr[^>]*>.*?<a:latin\b", xml, re.DOTALL))
            if latin:
                issues.append(f"C3: {name} has {latin} font-name overrides")
    return AssertionResult("C3", not issues, issues)


def check_c5_no_pictures_or_media(pptx_path: str) -> AssertionResult:
    """C5: Output must contain no picture shapes, blips, or embedded media."""
    issues: list[str] = []
    with zipfile.ZipFile(pptx_path) as zf:
        names = zf.namelist()
        media = [n for n in names if n.startswith("ppt/media/")]
        if media:
            issues.append(f"C5: package contains {len(media)} ppt/media files")

        xml_names = [
            n for n in names
            if n.startswith("ppt/") and n.endswith((".xml", ".rels"))
        ]
        for name in xml_names:
            xml = zf.read(name).decode("utf-8", errors="ignore")
            pic_count = xml.count("<p:pic")
            blip_count = xml.count("<a:blip")
            image_rel_count = xml.count("/image")
            if pic_count:
                issues.append(f"C5: {name} contains {pic_count} picture nodes")
            if blip_count:
                issues.append(f"C5: {name} contains {blip_count} image blips")
            if image_rel_count:
                issues.append(f"C5: {name} contains {image_rel_count} image relationships")
    return AssertionResult("C5", not issues, issues)


def check_c6_widescreen_16x9(pptx_path: str) -> AssertionResult:
    """C6: Output deck must be standard 16:9 widescreen."""
    prs = Presentation(pptx_path)
    width = int(prs.slide_width)
    height = int(prs.slide_height)
    ratio = width / height if height else 0
    issues: list[str] = []
    if abs(ratio - (16 / 9)) > 0.01:
        issues.append(f"C6: slide ratio {ratio:.3f} is not 16:9")
    return AssertionResult("C6", not issues, issues)


def _slide_fill_ratio(slide) -> float:
    """Very rough estimate: area occupied by shapes / slide area.

    Only counts visible shapes with width & height > 0. Good enough to detect
    "mostly-empty" decks.
    """
    slide_w = 13.333 * 914400
    slide_h = 7.5 * 914400
    total_area = slide_w * slide_h

    covered = 0
    for shape in slide.shapes:
        try:
            w = shape.width or 0
            h = shape.height or 0
        except Exception:
            continue
        covered += int(w) * int(h)
    return min(1.0, covered / total_area)


def check_c4_fill_ratio(pptx_path: str, lower: float = 0.30, upper: float = 1.4) -> AssertionResult:
    """C4: Body slides should neither be mostly empty nor overflowing."""
    prs = Presentation(pptx_path)
    issues: list[str] = []
    for i in range(1, len(prs.slides) - 1):
        ratio = _slide_fill_ratio(prs.slides[i])
        if ratio < lower:
            issues.append(f"C4: slide {i + 1} fill ratio {ratio:.2f} < {lower}")
        elif ratio > upper:
            issues.append(f"C4: slide {i + 1} fill ratio {ratio:.2f} > {upper}")
    return AssertionResult("C4", not issues, issues)


def run_all(pptx_path: str) -> dict[str, AssertionResult]:
    return {
        "C1": check_c1_fixed_slides(pptx_path),
        "C2": check_c2_no_adjacent_repeat(pptx_path),
        "C3": check_c3_master_inheritance(pptx_path),
        "C4": check_c4_fill_ratio(pptx_path),
        "C5": check_c5_no_pictures_or_media(pptx_path),
        "C6": check_c6_widescreen_16x9(pptx_path),
    }
