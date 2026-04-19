"""Main orchestrator: iterates PresentationContent.slides and dispatches to renderers.

This is the entry point for Step 4. Takes PresentationContent from Step 3 and
a template PPTX path, produces a .pptx file.

Includes infographic-first approach (3A): before rendering bullet slides,
checks if content can be visualized as an infographic instead.
"""

from pathlib import Path

from step2.slide_plan_models import SlideType, LayoutType
from step3.content_models import PresentationContent, SlideContent
from step4.template_manager import load_template, TemplateType, LayoutRole, get_layout
from step4.grid import Grid
from step4.renderers import (
    render_cover,
    render_end,
    render_bullets,
    render_chart,
    render_table,
    render_infographic,
    render_agenda,
    _add_slide_title,
    _add_key_message_footer,
    _get_bullet_texts,
)
from step4.infographic_renderers import (
    render_timeline,
    render_comparison_columns,
    render_numbered_steps,
    render_kpi_cards,
    render_funnel,
    render_pyramid,
    render_hub_spoke,
    render_cycle_diagram,
    render_vertical_list,
    render_pros_cons,
    render_matrix_2x2,
    render_icon_grid,
)
from step4.validator import validate_presentation
from step1.geo_detector import has_geographic_content

# Infographic-first: keywords that trigger infographic rendering
# instead of default bullet slides (research: hackathon_research.md:68-69)
INFOGRAPHIC_KEYWORDS: dict[str, str] = {
    "process": "numbered_steps",
    "steps": "numbered_steps",
    "workflow": "numbered_steps",
    "procedure": "numbered_steps",
    "how to": "numbered_steps",
    "methodology": "numbered_steps",
    "framework": "numbered_steps",
    "timeline": "timeline",
    "history": "timeline",
    "evolution": "timeline",
    "chronology": "timeline",
    "milestone": "timeline",
    "roadmap": "timeline",
    "phases": "timeline",
    "comparison": "comparison_columns",
    "vs": "comparison_columns",
    "versus": "comparison_columns",
    "contrast": "comparison_columns",
    "benchmark": "comparison_columns",
    "pros": "pros_cons",
    "cons": "pros_cons",
    "advantages": "pros_cons",
    "disadvantages": "pros_cons",
    "strengths": "pros_cons",
    "weaknesses": "pros_cons",
    "swot": "matrix_2x2",
    "hierarchy": "pyramid",
    "pyramid": "pyramid",
    "tier": "pyramid",
    "level": "pyramid",
    "kpi": "kpi_cards",
    "metric": "kpi_cards",
    "performance": "kpi_cards",
    "dashboard": "kpi_cards",
    "funnel": "funnel",
    "conversion": "funnel",
    "pipeline": "funnel",
    "cycle": "cycle_diagram",
    "circular": "cycle_diagram",
    "recurring": "cycle_diagram",
    "hub": "hub_spoke",
    "ecosystem": "hub_spoke",
    "components": "hub_spoke",
}


_ICON_KEYWORD_MAP: dict[str, str] = {
    "revenue": "finance", "profit": "finance", "cost": "finance", "budget": "finance",
    "growth": "growth", "increase": "growth", "rise": "growth", "expand": "growth",
    "decline": "decline", "decrease": "decline", "drop": "decline", "fall": "decline",
    "time": "time", "schedule": "time", "deadline": "time",
    "target": "target", "goal": "target", "objective": "target",
    "team": "people", "employee": "people", "staff": "people", "workforce": "people",
    "security": "shield", "protect": "shield", "safe": "shield",
    "global": "globe", "international": "globe", "world": "globe",
    "innovation": "lightbulb", "idea": "lightbulb", "creative": "lightbulb",
    "data": "database", "analytics": "chart_line", "metric": "chart_line",
    "cloud": "cloud", "infrastructure": "building", "technology": "code",
    "strategy": "target", "plan": "calendar", "launch": "rocket",
}


def _infer_icon_category(text: str) -> str:
    """Infer an icon category from bullet text using keyword matching.

    Includes geographic detection per judge feedback:
    "Flag could be added for countries" (Common Mistakes PPTX).
    """
    lower = text.lower()
    for keyword, category in _ICON_KEYWORD_MAP.items():
        if keyword in lower:
            return category
    if has_geographic_content(text):
        return "globe"
    return "check"


# Minimum bullet count required for each infographic type (Fix 12)
_MIN_ITEMS_FOR_INFOGRAPHIC: dict[str, int] = {
    "timeline": 2,
    "comparison_columns": 2,
    "numbered_steps": 2,
    "kpi_cards": 2,
    "funnel": 2,
    "pyramid": 2,
    "hub_spoke": 2,
    "cycle_diagram": 3,
    "vertical_list": 2,
    "pros_cons": 2,
    "matrix_2x2": 4,
    "icon_grid": 2,
}


def _detect_infographic_type(slide_content: SlideContent) -> str | None:
    """Infographic-first approach: check if content should be visualized.

    Returns the infographic type name or None for regular bullet rendering.
    Per research docs (hackathon_research.md:68-69): "Before placing text,
    ask: Can this be visualized?"

    Fix 12: Guards against triggering infographics when there aren't
    enough bullets to render them properly.
    """
    text = f"{slide_content.title} {slide_content.key_message}".lower()

    # Count available bullets before deciding
    bullets = getattr(slide_content, 'bullets', None)
    key_points = getattr(slide_content, 'key_points', None)
    bullet_count = len(bullets) if bullets else 0
    if key_points:
        bullet_count = max(bullet_count, len(key_points))

    for keyword, infographic_type in INFOGRAPHIC_KEYWORDS.items():
        if keyword in text:
            # Only enforce minimum if we can actually count bullets
            if bullet_count > 0:
                min_items = _MIN_ITEMS_FOR_INFOGRAPHIC.get(infographic_type, 2)
                if bullet_count < min_items:
                    return None
            return infographic_type

    return None


def build_presentation(
    content: PresentationContent,
    template_path: str,
    output_path: str,
) -> tuple[str, list[str]]:
    """Build a PPTX file from PresentationContent and a Slide Master template.

    Args:
        content: Structured slide content from Step 3.
        template_path: Path to one of the 3 Slide Master PPTX templates.
        output_path: Where to save the generated .pptx file.

    Returns:
        Tuple of (output_path, validation_issues).
    """
    prs, template_type = load_template(template_path)
    grid = Grid(template_type)

    print(f"  Template: {template_type.value}")
    print(f"  Slides to render: {len(content.slides)}")

    render_errors: list[str] = []
    for slide_content in content.slides:
        errors = _render_slide(prs, template_type, grid, slide_content, content)
        render_errors.extend(errors)

    if render_errors:
        print(f"\n  Render errors ({len(render_errors)}):")
        for err in render_errors:
            print(f"    - {err}")

    issues = validate_presentation(prs)
    issues.extend(render_errors)
    if issues:
        print(f"\n  Validation issues ({len(issues)}):")
        for issue in issues:
            print(f"    - {issue}")

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"\n  Saved: {output_path}")

    return output_path, issues


def _add_speaker_notes(prs, slide_content: SlideContent) -> None:
    """Add speaker notes to the most recently added slide.

    Uses key_message as the primary note, with title as context.
    gap_analysis.md listed speaker notes as "well addressed" — this implements it.
    """
    if not prs.slides:
        return
    slide = prs.slides[-1]
    notes_parts: list[str] = []
    if slide_content.key_message:
        notes_parts.append(slide_content.key_message)
    if slide_content.title and slide_content.title not in (slide_content.key_message or ""):
        notes_parts.append(f"Slide: {slide_content.title}")
    if not notes_parts:
        return
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        if tf is not None:
            tf.text = "\n".join(notes_parts)
    except Exception:
        pass


def _render_slide(
    prs,
    template_type: TemplateType,
    grid: Grid,
    slide_content: SlideContent,
    full_content: PresentationContent,
) -> list[str]:
    """Dispatch a single slide to the appropriate renderer.

    Includes infographic-first approach: before rendering bullets,
    checks if the content can be visualized as an infographic.

    VIZ-1: wraps each slide render in error protection.
    Adds speaker notes to every rendered slide.
    """
    errors: list[str] = []
    slide_num = slide_content.slide_number

    try:
        slide_type = slide_content.slide_type
        layout = slide_content.layout

        if slide_type == SlideType.TITLE:
            render_cover(prs, template_type, slide_content, full_content.title)
            _add_speaker_notes(prs, slide_content)
            return errors

        if slide_type == SlideType.THANK_YOU:
            render_end(prs, template_type, slide_content)
            _add_speaker_notes(prs, slide_content)
            return errors

        if slide_type == SlideType.AGENDA:
            render_agenda(prs, template_type, slide_content, grid, full_content.slides)
            _add_speaker_notes(prs, slide_content)
            return errors

        # Chart slide — has chart_data
        if slide_content.chart_data is not None:
            render_chart(prs, template_type, slide_content, grid)
            _add_speaker_notes(prs, slide_content)
            return errors

        # Table slide — has table_data but no chart
        if slide_content.table_data is not None:
            render_table(prs, template_type, slide_content, grid)
            _add_speaker_notes(prs, slide_content)
            return errors

        # Infographic layouts (comparison, timeline, process) — explicit layout
        if layout in (LayoutType.COMPARISON, LayoutType.TIMELINE, LayoutType.PROCESS):
            render_infographic(prs, template_type, slide_content, grid)
            _add_speaker_notes(prs, slide_content)
            return errors

        # Infographic-first approach: check if bullet content can be visualized
        infographic_type = _detect_infographic_type(slide_content)
        if infographic_type:
            _render_infographic_first(
                prs, template_type, slide_content, grid, infographic_type
            )
            _add_speaker_notes(prs, slide_content)
            return errors

        # Default: bullet slide
        render_bullets(prs, template_type, slide_content, grid)
        _add_speaker_notes(prs, slide_content)

    except Exception as exc:
        error_msg = f"Slide {slide_num} ({slide_content.slide_type.value}): render failed — {exc}"
        print(f"  WARNING: {error_msg}")
        errors.append(error_msg)

    return errors


def _render_infographic_first(
    prs,
    template_type: TemplateType,
    slide_content: SlideContent,
    grid: Grid,
    infographic_type: str,
) -> None:
    """Render a slide using an infographic renderer instead of plain bullets.

    This implements the infographic-first approach from the research docs:
    "Before placing text, ask: Can this be visualized?"
    """
    layout_obj = get_layout(prs, template_type, LayoutRole.CONTENT)
    slide = prs.slides.add_slide(layout_obj)

    _add_slide_title(slide, grid, template_type, slide_content.title)

    bullets = _get_bullet_texts(slide_content)

    # Build structured items from bullet text (Fix 7: richer data conversion)
    step_items = _bullets_to_steps(bullets)
    kpi_items = _bullets_to_kpis(bullets)
    timeline_items = _bullets_to_timeline(bullets)

    renderer_map = {
        "timeline": lambda: render_timeline(slide, grid, timeline_items),
        "comparison_columns": lambda: render_comparison_columns(
            slide, grid,
            [{"title": _split_title(b), "points": [_split_desc(b)]}
             for b in bullets]
        ),
        "numbered_steps": lambda: render_numbered_steps(slide, grid, step_items),
        "kpi_cards": lambda: render_kpi_cards(slide, grid, kpi_items),
        "funnel": lambda: render_funnel(
            slide, grid,
            [{"label": _split_title(b), "value": _extract_number(b)}
             for b in bullets]
        ),
        "pyramid": lambda: render_pyramid(slide, grid, bullets),
        "hub_spoke": lambda: render_hub_spoke(
            slide, grid, bullets[0] if bullets else "Center",
            bullets[1:] if len(bullets) > 1 else bullets
        ),
        "cycle_diagram": lambda: render_cycle_diagram(slide, grid, bullets),
        "vertical_list": lambda: render_vertical_list(slide, grid, step_items),
        "pros_cons": lambda: render_pros_cons(
            slide, grid,
            bullets[:len(bullets)//2],
            bullets[len(bullets)//2:]
        ),
        "matrix_2x2": lambda: render_matrix_2x2(
            slide, grid,
            [{"title": b, "points": []} for b in bullets[:4]]
        ),
        "icon_grid": lambda: render_icon_grid(
            slide, grid,
            [{"title": b, "description": "",
              "category": _infer_icon_category(b)} for b in bullets]
        ),
    }

    renderer = renderer_map.get(infographic_type)
    if renderer:
        renderer()
    else:
        # Fallback: render as regular bullets on the already-created slide
        from step4.font_sizing import auto_font_size_pt
        from pptx.util import Pt
        txBox = slide.shapes.add_textbox(
            grid.content_left, grid.content_top,
            grid.content_width, grid.content_height,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        all_text = " ".join(bullets)
        font_size = auto_font_size_pt(all_text, "full_width")
        for i, bullet_text in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {bullet_text}"
            p.font.size = font_size
            p.space_after = Pt(8)

    _add_key_message_footer(slide, grid, slide_content.key_message)


# ── Bullet-to-structured-data helpers (Fix 7) ────────────────────────

import re as _re


def _split_title(text: str) -> str:
    """Extract the title portion from a bullet (before colon/dash/period)."""
    for sep in (":", " – ", " — ", " - "):
        if sep in text:
            return text.split(sep, 1)[0].strip()
    words = text.split()
    return " ".join(words[:5]) if len(words) > 5 else text


def _split_desc(text: str) -> str:
    """Extract the description portion from a bullet (after colon/dash)."""
    for sep in (":", " – ", " — ", " - "):
        if sep in text:
            return text.split(sep, 1)[1].strip()
    return text


def _extract_number(text: str) -> str:
    """Extract the first number (with optional %, $, unit) from text."""
    m = _re.search(r'[\$]?[\d,]+\.?\d*\s*[%$BMKbmk]?', text)
    return m.group(0).strip() if m else ""


def _bullets_to_steps(bullets: list[str]) -> list[dict]:
    """Convert bullet texts to numbered-step items with title + description."""
    return [
        {"title": _split_title(b), "description": _split_desc(b)}
        for b in bullets
    ]


def _bullets_to_kpis(bullets: list[str]) -> list[dict]:
    """Convert bullet texts to KPI card items with value + label."""
    items = []
    for b in bullets:
        value = _extract_number(b)
        label = _split_title(b) if value else b
        if not value:
            value = "—"
        items.append({"value": value, "label": label})
    return items


def _bullets_to_timeline(bullets: list[str]) -> list[dict]:
    """Convert bullet texts to timeline items with label + description."""
    return [
        {"label": _split_title(b), "description": _split_desc(b)}
        for b in bullets
    ]
