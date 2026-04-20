"""Deck builder — single entry point that produces the 15-slide .pptx.

Pipeline:
    1. Load the template and drop its demo slides.
    2. Insert slide 1 via the template's cover layout; fill title/subtitle and
       add presenter + date at bottom-left (constraint C1).
    3. Schedule body slides 2..14 using the layout catalog with no-adjacent-
       same-class enforcement (constraint C2).
    4. Render each body slide via the chosen layout renderer — all styling
       cascades from the master (constraint C3).
    5. Append the template's thank-you layout as slide 15 (constraint C1).
    6. Persist the .pptx file.
"""

from __future__ import annotations

import json
import os
from typing import Optional

from pptx.enum.dml import MSO_THEME_COLOR

from step2.slide_plan_models import SlideType
from step3.content_models import PresentationContent, SlideContent

from step4 import layouts, template_ops
from step4.scheduler import schedule
from step4.template_ops import TemplateType


# Per-template primary accent — templates differ on which accent slot holds
# the dark brand color. UAE_Solar's accent1 is a near-white sage (#EFF3E5),
# so tiles would be illegible; its dark brand green lives in accent2.
_PRIMARY_ACCENT_BY_TEMPLATE: dict[TemplateType, MSO_THEME_COLOR] = {
    TemplateType.ACCENTURE: MSO_THEME_COLOR.ACCENT_1,
    TemplateType.AI_BUBBLE: MSO_THEME_COLOR.ACCENT_1,
    TemplateType.UAE_SOLAR: MSO_THEME_COLOR.ACCENT_2,
}


TOTAL_SLIDES = 15


def manifest_path_for(pptx_path: str) -> str:
    return pptx_path + ".layouts.json"


def _body_slices(content: PresentationContent) -> list[SlideContent]:
    """Return the 13 body slides (positions 2..14)."""
    ordered = sorted(content.slides, key=lambda s: s.slide_number)
    if len(ordered) < 3:
        # Degenerate case: pad via repetition so scheduler still gets 13 slots.
        pad = [ordered[-1]] * (TOTAL_SLIDES - len(ordered))
        ordered = ordered + pad

    # Strip any slide that is clearly a title or thank-you — those positions
    # are reserved for template slides.
    trimmed: list[SlideContent] = []
    for slide in ordered:
        if slide.slide_number in (1, TOTAL_SLIDES):
            continue
        if slide.slide_type in (SlideType.TITLE, SlideType.THANK_YOU):
            continue
        trimmed.append(slide)

    target = TOTAL_SLIDES - 2
    if len(trimmed) >= target:
        return trimmed[:target]

    # Pad by repeating the strongest slides so we always ship 13 body slides.
    if not trimmed:
        raise ValueError("No body slides produced by the designer")
    padded = list(trimmed)
    while len(padded) < target:
        padded.append(trimmed[len(padded) % len(trimmed)])
    return padded


def build_deck(
    content: PresentationContent,
    template_path: str,
    output_path: str,
    presenter: str,
    presentation_date: str,
    cover_subtitle: Optional[str] = None,
) -> str:
    """Produce the final 15-slide .pptx at ``output_path``."""
    prs, tpl = template_ops.load_blank_canvas(template_path)
    layouts.set_primary_accent(
        _PRIMARY_ACCENT_BY_TEMPLATE.get(tpl, MSO_THEME_COLOR.ACCENT_1)
    )

    # Slide 1 — cover
    subtitle = cover_subtitle or _derive_cover_subtitle(content)
    template_ops.add_cover_slide(
        prs,
        template=tpl,
        title=content.title,
        subtitle=subtitle,
        presenter=presenter,
        presentation_date=presentation_date,
    )

    # Slides 2..14 — scheduled body layouts
    body = _body_slices(content)
    assignments = schedule(body)
    layout_trace: list[dict] = [{"slide": 1, "layout": "cover", "class": "cover"}]
    for i, (slide_content, entry) in enumerate(zip(body, assignments)):
        # Alternate between the two "content" layouts for extra title-bar
        # variety when the master offers more than one.
        slide = template_ops.add_content_slide(prs, tpl, pick=i % 2)
        entry.render(slide, slide_content)
        layout_trace.append(
            {"slide": i + 2, "layout": entry.name, "class": entry.klass}
        )

    # Slide 15 — end
    template_ops.add_end_slide(prs, tpl)
    layout_trace.append({"slide": TOTAL_SLIDES, "layout": "end", "class": "end"})

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)

    with open(manifest_path_for(output_path), "w") as f:
        json.dump({"layouts": layout_trace, "template": tpl.value}, f, indent=2)
    return output_path


def _derive_cover_subtitle(content: PresentationContent) -> str:
    """Use the first body slide's subtitle or key message as the cover subtitle."""
    for slide in content.slides:
        if slide.subtitle:
            return slide.subtitle
    for slide in content.slides:
        if slide.key_message:
            return slide.key_message
    return ""
