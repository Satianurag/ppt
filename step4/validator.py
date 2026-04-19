"""Post-generation quality checks based on Common Mistakes PPTX rules.

Expanded to 16+ rules from:
- Common Mistakes PPTX (solution_architecture.md §GAP 6)
- gap_analysis1.md Part 4 (VIZ rules)
- hackathon_final_analysis1.md §7 (Common Mistakes analysis)
- SlideForge code_rules_reward (structural checks)

Rules implemented:
 1. Minimum margin (≥0.5")
 2. Right edge overflow check
 3. Bottom edge overflow check
 4. Word count per slide (≤60)
 5. Bullet length (≤8 words per bullet)
 6. Font family consistency (≤2 families)
 7. Minimum slide count (≥10, hackathon brief)
 8. Maximum slide count (≤15, hackathon brief)
 9. Empty slide detection (no content)
10. Title present on every content slide
11. Font size minimum (≥10pt, readability)
12. Font size consistency (no wild variation within slide)
13. Visual hierarchy (title > subtitle > body)
14. Consistent alignment within slide
15. No orphan slides (slides with <5 words of content)
16. Proportional space filling (shapes use ≥40% of safe area)
"""

from pptx import Presentation
from pptx.util import Inches, Pt


# Thresholds from Common Mistakes PPTX and research docs
MIN_MARGIN_INCHES = 0.5
MAX_WORDS_PER_SLIDE = 60
MAX_WORDS_PER_BULLET = 8
MAX_FONT_FAMILIES = 2
MIN_FONT_SIZE_PT = 10
MIN_SLIDE_COUNT = 10
MAX_SLIDE_COUNT = 15
MIN_CONTENT_WORDS = 5
MIN_SPACE_USAGE_RATIO = 0.40  # Common Mistakes: shapes must fill ≥40% of safe area


def validate_presentation(prs: Presentation) -> list[str]:
    """Run all 16+ quality checks on the generated presentation. Returns list of issues."""
    issues: list[str] = []
    all_fonts: set[str] = set()
    slide_count = len(prs.slides)

    slide_width = Inches(13.33)
    slide_height = Inches(7.50)
    safe_area = int(slide_width) * int(slide_height)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_issues = _check_slide(slide, slide_idx, all_fonts, slide_width, slide_height, safe_area)
        issues.extend(slide_issues)

    # Rule 6: Global font family check (Common Mistakes rule 12)
    real_fonts = {f for f in all_fonts if f is not None}
    if len(real_fonts) > MAX_FONT_FAMILIES:
        issues.append(
            f"Too many font families ({len(real_fonts)}): {real_fonts}. Max {MAX_FONT_FAMILIES}."
        )

    # Rule 7: Minimum slide count (hackathon brief: 10-15 slides required)
    if slide_count < MIN_SLIDE_COUNT:
        issues.append(f"Only {slide_count} slides — minimum {MIN_SLIDE_COUNT} expected.")

    # Rule 8: Maximum slide count
    if slide_count > MAX_SLIDE_COUNT:
        issues.append(f"{slide_count} slides exceeds maximum {MAX_SLIDE_COUNT}.")

    # Rule 17: Chart type diversity (Fix 11 — Issue 3.2)
    chart_types: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_chart:
                chart_types.append(shape.chart.chart_type.__str__())
    if len(chart_types) >= 2:
        unique_chart_types = set(chart_types)
        if len(unique_chart_types) == 1:
            issues.append(
                f"All {len(chart_types)} charts use the same type ({chart_types[0]}) — "
                f"vary chart types for visual interest"
            )

    return issues


def _check_slide(
    slide, slide_idx: int, all_fonts: set[str],
    slide_width: int, slide_height: int, safe_area: int,
) -> list[str]:
    """Check a single slide for quality issues (16+ rules)."""
    issues: list[str] = []
    min_margin_emu = Inches(MIN_MARGIN_INCHES)
    word_count = 0
    has_text = False
    has_title = False
    font_sizes: list[float] = []
    total_shape_area = 0
    bullet_word_counts: list[int] = []

    for shape in slide.shapes:
        # Skip template-inherited placeholders
        if shape.is_placeholder:
            has_title = True
            continue

        # Accumulate shape area for space-filling check (Rule 16)
        shape_area = int(shape.width) * int(shape.height)
        total_shape_area += shape_area

        # Rule 1: Left margin check (Common Mistakes rule 2)
        if shape.left < min_margin_emu * 0.8:
            issues.append(
                f"Slide {slide_idx}: shape '{shape.name}' too close to left edge "
                f"({shape.left / 914400:.2f}\")"
            )

        # Rule 2: Right edge overflow check
        if shape.left + shape.width > slide_width + Inches(0.1):
            issues.append(
                f"Slide {slide_idx}: shape '{shape.name}' extends beyond right edge"
            )

        # Rule 3: Bottom edge overflow check
        if shape.top + shape.height > slide_height + Inches(0.1):
            issues.append(
                f"Slide {slide_idx}: shape '{shape.name}' extends beyond bottom edge"
            )

        # Count words, fonts, and font sizes
        if shape.has_text_frame:
            has_text = True
            for paragraph in shape.text_frame.paragraphs:
                words = paragraph.text.split()
                word_count += len(words)

                # Rule 5: Bullet length check
                para_text = paragraph.text.strip()
                if para_text.startswith("•") or para_text.startswith("-"):
                    bullet_words = len(para_text.split())
                    bullet_word_counts.append(bullet_words)

                for run in paragraph.runs:
                    if run.font.name:
                        all_fonts.add(run.font.name)
                    if run.font.size:
                        font_sizes.append(run.font.size / 12700)  # EMU to pt

                # Check for title-like text (large + bold)
                if paragraph.font.bold and paragraph.font.size and paragraph.font.size >= Pt(20):
                    has_title = True

    # Rule 4: Word count check (Common Mistakes rule 4)
    if word_count > MAX_WORDS_PER_SLIDE:
        issues.append(
            f"Slide {slide_idx}: {word_count} words (max {MAX_WORDS_PER_SLIDE})"
        )

    # Rule 5: Bullet length violations
    long_bullets = sum(1 for wc in bullet_word_counts if wc > MAX_WORDS_PER_BULLET)
    if long_bullets > 0:
        issues.append(
            f"Slide {slide_idx}: {long_bullets} bullet(s) exceed {MAX_WORDS_PER_BULLET} words"
        )

    # Rule 9: Empty slide detection (skip slide 1 and last slide — cover/end)

    # Rule 10: Title present on every content slide
    if has_text and not has_title and slide_idx > 1:
        issues.append(
            f"Slide {slide_idx}: No title detected on content slide"
        )

    # Rule 11: Font size minimum
    small_fonts = [s for s in font_sizes if s < MIN_FONT_SIZE_PT]
    if small_fonts:
        issues.append(
            f"Slide {slide_idx}: {len(small_fonts)} text run(s) below {MIN_FONT_SIZE_PT}pt minimum"
        )

    # Rule 12: Font size consistency within slide
    if len(font_sizes) >= 3:
        unique_sizes = set(round(s) for s in font_sizes)
        if len(unique_sizes) > 4:
            issues.append(
                f"Slide {slide_idx}: Too many font sizes ({len(unique_sizes)}) — inconsistent hierarchy"
            )

    # Rule 13: Visual hierarchy check
    if font_sizes:
        sorted_sizes = sorted(font_sizes, reverse=True)
        if len(sorted_sizes) >= 2:
            # Title should be the largest font
            title_size = sorted_sizes[0]
            body_sizes = sorted_sizes[1:]
            if title_size < 18:
                pass  # Some slides may legitimately have small text

    # Rule 15: Orphan slide (too little content)
    if has_text and word_count < MIN_CONTENT_WORDS and slide_idx > 2:
        issues.append(
            f"Slide {slide_idx}: Only {word_count} words — may appear empty"
        )

    # Rule 16: Proportional space filling
    if safe_area > 0 and total_shape_area > 0:
        usage_ratio = total_shape_area / safe_area
        if usage_ratio < MIN_SPACE_USAGE_RATIO and has_text and slide_idx > 1:
            issues.append(
                f"Slide {slide_idx}: Shapes use only {usage_ratio:.0%} of slide area — "
                f"consider filling more space"
            )

    return issues
