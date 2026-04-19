"""Reviewer Agent — quality assurance using SlideForge scoring.

Wraps the quality scoring and validation systems as a dedicated agent role.
Inspired by SlideForge's 6-component reward system and PPTAgent's retry pattern.

Responsibilities:
  - Score the rendered presentation using SlideForge 6-component system
  - Validate against Common Mistakes rules (16+ checks)
  - Generate quality report
  - Decide pass/fail and provide feedback for retry
"""

from agents.base import BaseAgent
from agents.protocol import AgentRole, MessageType, PipelineState
from step3.content_optimizer import score_presentation
from step4.validator import validate_presentation

from pptx import Presentation


class ReviewerAgent(BaseAgent):
    """Reviews presentation quality and decides pass/fail.

    Uses SlideForge's 6-component scoring system:
      structural_rules (1.0), content_quality (2.0), render_quality (2.0),
      brief_reconstruction (2.0), source_coverage (1.5), narrative_flow (1.0)

    Combined with 16-rule Common Mistakes validator.
    """

    PASS_THRESHOLD = 0.6

    def __init__(self, threshold: float = 0.6) -> None:
        super().__init__(role=AgentRole.REVIEWER, name="Reviewer")
        self.threshold = threshold

    def process(self, state: PipelineState) -> PipelineState:
        """Review the rendered PPTX for quality.

        Reads: state.pptx_path, state.presentation_content, state.render_issues
        Writes: state.quality_score, state.quality_report, state.review_passed, state.review_feedback
        """
        content = state.presentation_content
        pptx_path = state.pptx_path

        self.send_message(
            state, AgentRole.COORDINATOR, MessageType.STATUS,
            {"phase": "review", "pptx": pptx_path},
        )

        # 1. Validate rendered PPTX with Common Mistakes rules
        if pptx_path:
            prs = Presentation(pptx_path)
            validation_issues = validate_presentation(prs)
        else:
            validation_issues = ["No PPTX file to validate"]

        # 2. Score content quality using standalone scoring (no LLM needed)
        if content is not None:
            overall, scores, quality_report = score_presentation(content)
        else:
            overall, scores, quality_report = 0.0, [], "No content to score."

        state.quality_score = overall
        state.quality_report = quality_report

        # 3. Determine pass/fail. Critical = anything the hackathon judges
        # will visibly deduct for: overflow, margin violations, empty slides,
        # low space usage, layout monotony.
        # A validation issue is "critical" when it maps to a hackathon deduction
        # that the Designer can realistically fix via re-extraction. Single
        # empty-space warnings are noisy (retry cannot meaningfully increase
        # content on one slide) — only treat them as critical when systemic.
        hard_tokens = ("overflow", "margin violation", "no title", "orphan",
                       "empty slide", "no text")
        fill_warnings = [i for i in validation_issues
                         if "fill more" in i.lower()
                         or "space usage" in i.lower()
                         or "consider filling" in i.lower()]
        critical_issues = [i for i in validation_issues
                           if any(tok in i.lower() for tok in hard_tokens)]
        # Systemic space-fill problem: 3+ slides under-filled at once.
        if len(fill_warnings) >= 3:
            critical_issues.append(
                f"{len(fill_warnings)} slides under-fill the canvas — "
                "add more bullets or use an infographic layout"
            )

        # Layout diversity check — hackathon Feedback #2 hard gate.
        unique_layouts: set[str] = set()
        if pptx_path:
            for sl in prs.slides:
                unique_layouts.add(sl.slide_layout.name or "")
        if len(unique_layouts) < 4 and len(prs.slides) >= 8:
            critical_issues.append(
                f"Only {len(unique_layouts)} unique layouts across "
                f"{len(prs.slides)} slides — need more variety"
            )

        passed = overall >= self.threshold and len(critical_issues) == 0
        state.review_passed = passed

        # 4. Generate targeted feedback for Designer retry. The feedback is
        # consumed verbatim by the extractor LLM prompt, so it must be
        # specific and actionable.
        feedback_parts: list[str] = []
        if not passed:
            if overall < self.threshold:
                feedback_parts.append(
                    f"Overall quality {overall:.2f} is below {self.threshold}."
                )
            if critical_issues:
                feedback_parts.append(
                    "Fix these specific slides: " + "; ".join(critical_issues[:6])
                )
            if scores:
                avg = lambda attr: sum(getattr(s, attr) for s in scores) / len(scores)
                low = [
                    f"{c}={avg(c):.2f}"
                    for c in ("structural_rules", "content_quality",
                              "render_quality", "source_coverage", "narrative_flow")
                    if avg(c) < 0.5
                ]
                if low:
                    feedback_parts.append("Weak dimensions: " + ", ".join(low))

        state.review_feedback = " | ".join(feedback_parts) if feedback_parts else "PASSED"

        self.record_turn(
            input_summary=f"Review {pptx_path}",
            output_summary=f"Score: {overall:.2f}, Passed: {passed}, Issues: {len(validation_issues)}",
            success=passed,
        )

        msg_type = MessageType.RESPONSE if passed else MessageType.FEEDBACK
        self.send_message(
            state, AgentRole.COORDINATOR, msg_type,
            {
                "score": overall,
                "passed": passed,
                "validation_issues": len(validation_issues),
                "critical_issues": len(critical_issues),
                "feedback": state.review_feedback,
            },
        )

        return state
